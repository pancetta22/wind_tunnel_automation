#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""260608 実験結果まとめ（簡潔版）スライド生成"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

C_DARK=RGBColor(0x1A,0x37,0x6C); C_MID=RGBColor(0x2E,0x6E,0xB4)
C_GREEN=RGBColor(0x21,0x7A,0x3C); C_GREEN_L=RGBColor(0xD6,0xEE,0xDF)
C_ORANGE=RGBColor(0xD9,0x6D,0x00); C_ORANGE_L=RGBColor(0xFD,0xF0,0xDE)
C_GRAY=RGBColor(0xF4,0xF6,0xF9); C_TEXT=RGBColor(0x1A,0x1A,0x2E)
C_WHITE=RGBColor(0xFF,0xFF,0xFF); C_RED=RGBColor(0xC0,0x39,0x2B); C_SUB=RGBColor(0x5A,0x6A,0x85)

prs=Presentation(); prs.slide_width=Inches(13.33); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]

def rect(s,l,t,w,h,fill=None,border=None,bw=Pt(1)):
    sh=s.shapes.add_shape(1,l,t,w,h); sh.shadow.inherit=False
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    else: sh.fill.background()
    if border: sh.line.color.rgb=border; sh.line.width=bw
    else: sh.line.fill.background()
    return sh

def text(s,t,l,tp,w,h,size=Pt(14),bold=False,color=C_TEXT,align=PP_ALIGN.LEFT,it=False):
    tb=s.shapes.add_textbox(l,tp,w,h); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=t
    r.font.size=size; r.font.bold=bold; r.font.color.rgb=color; r.font.italic=it
    return tb

def header(s,title,sub=None):
    rect(s,0,Inches(0.85),prs.slide_width,prs.slide_height-Inches(0.85),fill=C_GRAY)
    rect(s,0,0,prs.slide_width,Inches(0.85),fill=C_DARK)
    text(s,title,Inches(0.35),Inches(0.08),Inches(12.5),Inches(0.55),size=Pt(23),bold=True,color=C_WHITE)
    if sub: text(s,sub,Inches(0.37),Inches(0.56),Inches(12),Inches(0.3),size=Pt(11.5),color=RGBColor(0xA8,0xC8,0xF0))

def bullets(s,lines,l,t,w,h,size=Pt(13),gap=Pt(6)):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True; first=True
    for it in lines:
        txt=it[0]; lv=it[1] if len(it)>1 else 0; col=it[2] if len(it)>2 else C_TEXT; bd=it[3] if len(it)>3 else False
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.space_before=gap; p.level=lv
        mark="■ " if lv==0 else "・ "
        r=p.add_run(); r.text=mark+txt; r.font.size=size if lv==0 else Pt(size.pt-1.5)
        r.font.color.rgb=col; r.font.bold=bd

def panel(s,l,t,w,h,fill,border): rect(s,l,t,w,h,fill=fill,border=border,bw=Pt(1.5))

# ===== S1 サマリ =====
s=prs.slides.add_slide(BLANK)
header(s,"260608 実験結果：過去データへ大きく改善","風速も旧レベルに復帰。揚力傾斜 1.29× → 1.14×")

panel(s,Inches(0.4),Inches(1.1),Inches(4.0),Inches(2.5),C_GREEN_L,C_GREEN)
text(s,"揚力傾斜 dCl/dα",Inches(0.55),Inches(1.2),Inches(3.7),Inches(0.4),size=Pt(13),bold=True,color=C_GREEN)
text(s,"4.30",Inches(0.55),Inches(1.65),Inches(3.7),Inches(0.9),size=Pt(44),bold=True,color=C_DARK)
text(s,"/rad（旧3.76 / 新605=4.86）",Inches(0.55),Inches(2.6),Inches(3.7),Inches(0.4),size=Pt(12),color=C_SUB)
text(s,"旧比 1.14×（605は1.29×）",Inches(0.55),Inches(2.95),Inches(3.7),Inches(0.4),size=Pt(13),bold=True,color=C_GREEN)

panel(s,Inches(4.6),Inches(1.1),Inches(4.0),Inches(2.5),C_GREEN_L,C_GREEN)
text(s,"ゼロ揚力角 α₀",Inches(4.75),Inches(1.2),Inches(3.7),Inches(0.4),size=Pt(13),bold=True,color=C_GREEN)
text(s,"0.93°",Inches(4.75),Inches(1.65),Inches(3.7),Inches(0.9),size=Pt(44),bold=True,color=C_DARK)
text(s,"（旧0.5〜0.7 / 新605=1.45）",Inches(4.75),Inches(2.6),Inches(3.7),Inches(0.4),size=Pt(12),color=C_SUB)
text(s,"NACA0012理論0°へ接近",Inches(4.75),Inches(2.95),Inches(3.7),Inches(0.4),size=Pt(13),bold=True,color=C_GREEN)

panel(s,Inches(8.8),Inches(1.1),Inches(4.1),Inches(2.5),C_GREEN_L,C_GREEN)
text(s,"最大揚力 Cl_max",Inches(8.95),Inches(1.2),Inches(3.8),Inches(0.4),size=Pt(13),bold=True,color=C_GREEN)
text(s,"0.914",Inches(8.95),Inches(1.65),Inches(3.8),Inches(0.9),size=Pt(44),bold=True,color=C_DARK)
text(s,"@16°（旧0.81 / 新605=1.045）",Inches(8.95),Inches(2.6),Inches(3.8),Inches(0.4),size=Pt(12),color=C_SUB)
text(s,"旧へ接近",Inches(8.95),Inches(2.95),Inches(3.8),Inches(0.4),size=Pt(13),bold=True,color=C_GREEN)

panel(s,Inches(0.4),Inches(3.8),Inches(12.5),Inches(3.05),C_WHITE,C_MID)
text(s,"ポイント",Inches(0.55),Inches(3.9),Inches(12),Inches(0.4),size=Pt(15),bold=True,color=C_DARK)
bullets(s,[
 ("605→608で全指標（傾斜・α₀・Cl_max）が一斉に旧へ接近",0,C_GREEN,True),
 ("風速も旧レベルに復帰（mV 1056→1171、U 11.65→12.29 m/s）",0),
 ("結果が「実験のやり直し（物理セットアップ変更）」で動いた → マウント依存＝物理要因の強い証拠",0,C_RED,True),
 ("ソフト・後処理は固定（不変）なので、変わったのは物理側のみ",0,C_SUB),
 ("ただし まだ旧比 1.14× 残存 → 完全解決ではない",0,C_ORANGE,True),
],Inches(0.55),Inches(4.35),Inches(12.1),Inches(2.4),size=Pt(13.5),gap=Pt(9))

# ===== S2 比較表 =====
s=prs.slides.add_slide(BLANK)
header(s,"過去データとの比較（剛体 NACA0012）")
data=[
 ("旧260520","1165","12.23","3.76","0.74","0.81","—",False),
 ("旧260430","1166","12.08","3.76","0.69","0.81","—",False),
 ("旧251020","1163","12.11","3.77","0.52","0.79","—",False),
 ("新605","1056","11.65","4.86","1.45","1.045","1.29×",False),
 ("新605_2","1041","11.55","4.89","1.36","1.058","1.30×",False),
 ("新608 ★","1171","12.29","4.30","0.93","0.914","1.14×",True),
]
heads=["実験","mV","U[m/s]","dCl/dα[/rad]","α₀[°]","Cl_max","旧比"]
xs=[0.5,2.6,3.7,5.0,7.3,8.6,10.0]; ws=[2.1,1.1,1.3,2.3,1.3,1.4,1.6]
y=Inches(1.2)
rect(s,Inches(0.4),y,Inches(12.5),Inches(0.42),fill=C_DARK)
for hx,hw,ht in zip(xs,ws,heads):
    text(s,ht,Inches(hx),y+Inches(0.04),Inches(hw),Inches(0.34),size=Pt(12.5),bold=True,color=C_WHITE)
y=Inches(1.66)
for i,(lab,mv,U,rad,a0,clm,ratio,hl) in enumerate(data):
    bg=C_ORANGE_L if hl else (C_WHITE if i%2==0 else RGBColor(0xEC,0xF1,0xF8))
    rect(s,Inches(0.4),y,Inches(12.5),Inches(0.62),fill=bg,border=(C_ORANGE if hl else RGBColor(0xD0,0xDA,0xE8)),bw=Pt(2 if hl else 0.75))
    col=C_RED if hl else C_TEXT
    vals=[lab,mv,U,rad,a0,clm,ratio]
    for hx,hw,v in zip(xs,ws,vals):
        text(s,v,Inches(hx),y+Inches(0.13),Inches(hw),Inches(0.4),size=Pt(13),bold=hl,color=col)
    y+=Inches(0.62)

panel(s,Inches(0.4),Inches(5.7),Inches(12.5),Inches(1.15),C_GREEN_L,C_GREEN)
text(s,"傾向：旧(3.76) ← 新608(4.30) ← 新605(4.86)。608は旧と605の中間で、旧側へ大きく移動。風速も旧と一致。",
     Inches(0.6),Inches(5.85),Inches(12.1),Inches(0.5),size=Pt(13.5),bold=True,color=C_DARK)
text(s,"※ 260608はPofst_00.00重複（実験再開時の余りファイル）を除外して処理",
     Inches(0.6),Inches(6.35),Inches(12),Inches(0.4),size=Pt(11),it=True,color=C_SUB)

# ===== S3 解釈と次の一手 =====
s=prs.slides.add_slide(BLANK)
header(s,"解釈と次の一手")
panel(s,Inches(0.4),Inches(1.1),Inches(6.15),Inches(3.4),C_WHITE,C_MID)
text(s,"何がわかったか",Inches(0.55),Inches(1.2),Inches(5.8),Inches(0.4),size=Pt(15),bold=True,color=C_DARK)
bullets(s,[
 ("再マウント（物理変更）で傾斜・α₀が同時に改善",0,C_GREEN,True),
 ("→ 揚力測定は「マウント／水平荷重経路」に依存",0),
 ("ソフト・処理・風速ではない（固定／正常）",0,C_SUB),
 ("これまでの容疑除外と完全に整合：",0),
 ("残る本命＝マウント幾何・水平荷重経路",1,C_RED),
],Inches(0.55),Inches(1.65),Inches(5.85),Inches(2.7),size=Pt(13))

panel(s,Inches(6.75),Inches(1.1),Inches(6.15),Inches(3.4),C_ORANGE_L,C_ORANGE)
text(s,"まだ残る課題",Inches(6.9),Inches(1.2),Inches(5.8),Inches(0.4),size=Pt(15),bold=True,color=C_ORANGE)
bullets(s,[
 ("旧比 1.14× がまだ残存（605の1.29×からは改善）",0,C_RED,True),
 ("Cl_max・α₀も旧に「近い」が一致はしていない",0),
 ("マウントの詰め（剛性・作用線・再現性）でさらに旧へ寄せられる可能性",0),
],Inches(6.9),Inches(1.65),Inches(5.85),Inches(2.7),size=Pt(13))

panel(s,Inches(0.4),Inches(4.7),Inches(12.5),Inches(2.1),RGBColor(0x14,0x2A,0x55),C_MID)
text(s,"次の一手",Inches(0.6),Inches(4.82),Inches(11),Inches(0.4),size=Pt(15),bold=True,color=C_WHITE)
bullets(s,[
 ("605→608で『何を変えたか』を記録（再マウント有無・治具・締結）",0,C_WHITE,True),
 ("Fy軸の清浄な荷重テスト（鉛直重力）でセンサゲインの最終確認",0,RGBColor(0xC8,0xF0,0xD8)),
 ("マウントの作用線・剛性・再現性を詰め、残り1.14×を旧へ寄せる",0,RGBColor(0xC8,0xF0,0xD8)),
],Inches(0.6),Inches(5.25),Inches(12.1),Inches(1.5),size=Pt(13.5),gap=Pt(8))

prs.save("Windy_260608結果.pptx")
print("生成完了: Windy_260608結果.pptx (%dスライド)"%len(prs.slides._sldIdLst))
