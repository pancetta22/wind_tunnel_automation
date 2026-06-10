#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
make_rigid_comparison_pptx.py
NACA0012 rigid wing 全データ比較 パワポ生成
Ito(240521/240603/241223) + 250924/251020 + 260417/260424/260430/260520 + 260605/260605_2
260605/260605_2 は風洞自動化新システムによる計測
"""

import os, io
import numpy as np
import pandas as pd
from scipy.interpolate import interp1d
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.font_manager as _fm
# 利用可能な日本語フォントを自動選択（Mac=Hiragino / Windows=Yu Gothic 等）
_avail = {f.name for f in _fm.fontManager.ttflist}
for _jp in ["Hiragino Sans", "Yu Gothic", "Meiryo", "MS Gothic",
            "Noto Sans CJK JP", "IPAexGothic", "TakaoGothic"]:
    if _jp in _avail:
        plt.rcParams["font.family"] = _jp
        break
import matplotlib.ticker as ticker
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn as _qn

# ─── 定数 ────────────────────────────────────────────────────────────────────
# --- 考察フォルダ内で自己完結（aero_data/ の C_aero.csv を参照、テンプレ・出力も同フォルダ）---
SCRIPT_DIR    = os.path.dirname(os.path.abspath(__file__))
BASE          = os.path.join(SCRIPT_DIR, "aero_data")   # 各実験 C_aero.csv の置き場
TEMPLATE_PATH = os.path.join(SCRIPT_DIR, "研究室MTGテンプレート.pptx")
OUT           = os.path.join(SCRIPT_DIR, "Windy_rigid3_研究室比較.pptx")

SLIDE_W = Inches(10.0)
SLIDE_H = Inches(7.5)

# ─── データ定義 ───────────────────────────────────────────────────────────────
DATASETS = {
    "Ito\n240521":  {"path": os.path.join(BASE, "force_measurement_240521_ito_rigid/C_aero.csv"),
                     "color": "#9370DB", "marker": "v",  "date": "2024-05-21", "who": "ito"},
    "Ito\n240603":  {"path": os.path.join(BASE, "force_measurement_240603_ito_rigid/C_aero.csv"),
                     "color": "#7B2FBE", "marker": "<",  "date": "2024-06-03", "who": "ito"},
    "Ito\n241223":  {"path": os.path.join(BASE, "force_measurement_241223_ito_rigid/C_aero.csv"),
                     "color": "#4B0082", "marker": ">",  "date": "2024-12-23", "who": "ito"},
    "250924":       {"path": os.path.join(BASE, "force_measurement_250924_rigid/C_aero.csv"),
                     "color": "#4169E1", "marker": "^",  "date": "2025-09-24", "who": "okamoto"},
    "251020":       {"path": os.path.join(BASE, "force_measurement_251020_rigid/C_aero.csv"),
                     "color": "#E84040", "marker": "o",  "date": "2025-10-20", "who": "okamoto"},
    "260417":       {"path": os.path.join(BASE, "force_measurement_260417/C_aero.csv"),
                     "color": "#2E8B57", "marker": "s",  "date": "2026-04-17", "who": "kaneko"},
    "260424":       {"path": os.path.join(BASE, "force_measurement_260424/C_aero.csv"),
                     "color": "#FF8C00", "marker": "D",  "date": "2026-04-24", "who": "kaneko"},
    "260430":       {"path": os.path.join(BASE, "force_measurement_260430_rigid/C_aero.csv"),
                     "color": "#DC143C", "marker": "P",  "date": "2026-04-30", "who": "okamoto"},
    "260520":       {"path": os.path.join(BASE, "force_measurement_260520_rigid/C_aero.csv"),
                     "color": "#FF1493", "marker": "h",  "date": "2026-05-20", "who": "kaneko"},
    "260605①":     {"path": os.path.join(BASE, "force_measurement_260605_rigid/C_aero.csv"),
                     "color": "#00BFFF", "marker": "X",  "date": "2026-06-05", "who": "new_sys"},
    "260605②":     {"path": os.path.join(BASE, "force_measurement_260605_rigid2/C_aero.csv"),
                     "color": "#00CED1", "marker": "*",  "date": "2026-06-05", "who": "new_sys"},
    "260608":       {"path": os.path.join(BASE, "force_measurement_260608_rigid/C_aero.csv"),
                     "color": "#FFA500", "marker": "p",  "date": "2026-06-08", "who": "new_sys"},
    "rigid3\n修正後": {"path": os.path.join(BASE, "force_measurement_260610_rigid3/C_aero.csv"),
                     "color": "#00A000", "marker": "*",  "date": "2026-06-10", "who": "fixed"},
}

DISP_NAMES = {
    "Ito\n240521":  "240521",
    "Ito\n240603":  "240603",
    "Ito\n241223":  "241223",
    "250924":       "250924",
    "251020":       "251020",
    "260417":       "260417",
    "260424":       "260424",
    "260430":       "260430",
    "260520":       "260520",
    "260605①":     "260605①",
    "260605②":     "260605②",
    "260608":       "260608",
    "rigid3\n修正後": "rigid3(修正)",
}

ROW_INFO = [
    ("伊東さん 240521", "ito"),
    ("伊東さん 240603", "ito"),
    ("伊東さん 241223", "ito"),
    ("岡本    250924",  "okamoto"),
    ("岡本    251020",  "okamoto"),
    ("金光    260417",  "kaneko"),
    ("金光    260424",  "kaneko"),
    ("岡本    260430",  "okamoto"),
    ("金光    260520",  "kaneko"),
    ("260605① 新Sys",  "new_sys"),
    ("260605② 新Sys",  "new_sys"),
    ("260608 新Sys",   "new_sys"),
    ("rigid3 修正後",  "fixed"),
]

BG_COLOR = {
    "ito":      RGBColor(0xEA, 0xEF, 0xFF),
    "kaneko":   RGBColor(0xEA, 0xFF, 0xEF),
    "okamoto":  RGBColor(0xFF, 0xF5, 0xEA),
    "new_sys":  RGBColor(0xE0, 0xF8, 0xFF),
    "fixed":    RGBColor(0xD6, 0xF5, 0xD6),
}
NOTE_STR = {
    "ito":      "（伊東さん）",
    "kaneko":   "（金光くん）",
    "okamoto":  "（岡本）",
    "new_sys":  "（新Sys・接触不良）",
    "fixed":    "（接触修正後）",
}

# ─── aero_data 内の未登録フォルダを自動追加 ──────────────────────────────────────
#  今後 update_aero_data.py で aero_data/ に新しい rigid 実験の C_aero.csv が
#  追加されると、ここで自動的に DATASETS へ取り込まれ、表・全図に反映される。
#  （既存の登録済み13件はそのままのスタイル・解説を維持）
import re as _re
_known_paths = {os.path.normpath(v["path"]) for v in DATASETS.values()}
_auto_palette = ["#666666", "#8B4513", "#6A0DAD", "#008B8B", "#B8860B",
                 "#2F4F4F", "#C71585", "#556B2F", "#483D8B", "#A0522D"]
_auto_markers = ["o", "s", "^", "D", "v", "p", "X", "h", "<", ">"]
_ai = 0
if os.path.isdir(BASE):
    for _sub in sorted(os.listdir(BASE)):
        _ca = os.path.join(BASE, _sub, "C_aero.csv")
        if not os.path.isfile(_ca):
            continue
        if os.path.normpath(_ca) in _known_paths:
            continue
        _m    = _re.search(r"(\d{6})", _sub)
        _date = f"20{_m.group(1)[:2]}-{_m.group(1)[2:4]}-{_m.group(1)[4:]}" if _m else "—"
        DATASETS[_sub]   = {"path": _ca,
                            "color":  _auto_palette[_ai % len(_auto_palette)],
                            "marker": _auto_markers[_ai % len(_auto_markers)],
                            "date":   _date, "who": "new"}
        _disp            = _sub.replace("force_measurement_", "").replace("_rigid", "")
        DISP_NAMES[_sub] = _disp
        ROW_INFO.append((_disp + " 新規", "new"))
        _ai += 1
BG_COLOR.setdefault("new", RGBColor(0xEC, 0xEC, 0xEC))
NOTE_STR.setdefault("new", "（新規・自動追加）")
if _ai:
    print(f"[自動追加] aero_data 内の未登録 {_ai} 件を取り込みました。")

# ─── データ読み込み ───────────────────────────────────────────────────────────
def load(info):
    df = pd.read_csv(info["path"])
    df = df.apply(pd.to_numeric, errors="coerce").dropna()
    return df

data = {k: load(v) for k, v in DATASETS.items()}

# ─── 統計量計算 ───────────────────────────────────────────────────────────────
stats = {}
for k, df in data.items():
    row0  = df[df["AoA"] == 0]
    cl0   = float(row0.iloc[0]["Cl"]) if len(row0) > 0 else np.nan
    cd0   = float(row0.iloc[0]["Cd"]) if len(row0) > 0 else np.nan
    cm0   = float(row0.iloc[0]["Cm"]) if len(row0) > 0 else np.nan
    seg   = df[(df["AoA"] >= -5) & (df["AoA"] <= 5)].sort_values("AoA_mod")
    slope = float(np.polyfit(seg["AoA_mod"].astype(float),
                              seg["Cl"].astype(float), 1)[0]) if len(seg) >= 2 else np.nan
    # 零揚力迎角
    aoa0 = None
    for i in range(len(seg) - 1):
        c1, c2 = float(seg.iloc[i]["Cl"]), float(seg.iloc[i+1]["Cl"])
        if c1 <= 0 <= c2 or c2 <= 0 <= c1:
            am1 = float(seg.iloc[i]["AoA_mod"])
            am2 = float(seg.iloc[i+1]["AoA_mod"])
            aoa0 = am1 + (0 - c1) / (c2 - c1) * (am2 - am1)
            break
    if aoa0 is None:
        aoa0 = 0.0
    # P-M 対称性
    pos = df[df["AoA"] >= 0].sort_values("AoA")
    neg = df[df["AoA"] <= 0].sort_values("AoA")
    if len(neg) > 3:
        aoa_pm = np.arange(1, 21, dtype=float)
        fp = np.interp(aoa_pm, pos["AoA"].values, pos["Cl"].values,
                       left=np.nan, right=np.nan)
        fn = np.interp(aoa_pm, -neg["AoA"].values[::-1], neg["Cl"].values[::-1],
                       left=np.nan, right=np.nan)
        diff = fp - (-fn)
        mask = ~np.isnan(diff)
        pm_rmse = float(np.sqrt(np.mean(diff[mask]**2))) if mask.any() else np.nan
    else:
        pm_rmse = np.nan
    stats[k] = {"Cl0": cl0, "Cd0": cd0, "Cm0": cm0,
                 "aoa0": aoa0, "slope": slope, "pm_rmse": pm_rmse,
                 "date": DATASETS[k]["date"], "color": DATASETS[k]["color"]}

# ─── 零揚力迎角補正後 P-M RMSE ───────────────────────────────────────────────
pm_rmse_corr = {}
for k, df in data.items():
    df_s = df.sort_values("AoA_mod")
    a0 = stats[k]["aoa0"]
    deltas = np.arange(1, 21, dtype=float)
    f = interp1d(df_s["AoA_mod"].astype(float), df_s["Cl"].astype(float),
                 kind="linear", bounds_error=False, fill_value=np.nan)
    v_plus  = f(a0 + deltas)
    v_minus = f(a0 - deltas)
    diff = v_plus - (-v_minus)
    mask = ~(np.isnan(v_plus) | np.isnan(v_minus))
    pm_rmse_corr[k] = float(np.sqrt(np.mean(diff[mask]**2))) if mask.any() else np.nan

labels  = list(DATASETS.keys())
x_pos   = np.arange(len(labels))
colors  = [DATASETS[k]["color"] for k in labels]
disp    = [DISP_NAMES.get(k, k.replace("\n", " ")) for k in labels]

# ─── 図作成ユーティリティ ─────────────────────────────────────────────────────
def fig_to_bytes(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=180, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf

TS  = 12
LW  = 1.5
MK  = 4

# ─── 図1: 零揚力迎角の時系列 ─────────────────────────────────────────────────
fig, ax = plt.subplots(figsize=(10, 5))
vals = [stats[k]["aoa0"] for k in labels]
bar_colors = [DATASETS[k]["color"] for k in labels]
bars = ax.bar(x_pos, vals, color=bar_colors, edgecolor="white", linewidth=0.8, width=0.7)
for bar, v in zip(bars, vals):
    ax.text(bar.get_x() + bar.get_width()/2, v + 0.01,
            f"{v:.3f}°", ha="center", va="bottom", fontsize=9, fontweight="bold")
ax.axhline(0, color="black", linewidth=0.8, linestyle="--")
# 260605 背景強調
for i, k in enumerate(labels):
    if "260605" in k:
        ax.axvspan(i - 0.45, i + 0.45, alpha=0.12, color="deepskyblue", zorder=0)
ax.set_xticks(x_pos)
ax.set_xticklabels(disp, fontsize=TS - 2, rotation=15, ha="right")
ax.set_ylabel("零揚力迎角 [deg]", fontsize=TS)
ax.set_ylim(-0.1, 1.2)
ax.set_title("零揚力迎角の時系列変化（Cl=0 となる AoA_mod）", fontsize=TS)
ax.yaxis.set_major_locator(ticker.MultipleLocator(0.2))
ax.grid(axis="y", linestyle="--", alpha=0.5)
fig.tight_layout()
img_aoa0 = fig_to_bytes(fig)

# ─── 図2: Clスロープの時系列 ─────────────────────────────────────────────────
fig, ax = plt.subplots(figsize=(10, 5))
slopes = [stats[k]["slope"] for k in labels]
bars = ax.bar(x_pos, slopes, color=bar_colors, edgecolor="white", linewidth=0.8, width=0.7)
for bar, v in zip(bars, slopes):
    if not np.isnan(v):
        ax.text(bar.get_x() + bar.get_width()/2, v + 0.0005,
                f"{v:.4f}", ha="center", va="bottom", fontsize=9, fontweight="bold")
# 2π/rad = ~0.10966/deg 理論線
ax.axhline(2 * np.pi / 180, color="gray", linewidth=1.2, linestyle=":",
           label=r"$2\pi$ /rad = {:.4f}/deg".format(2*np.pi/180))
for i, k in enumerate(labels):
    if "260605" in k:
        ax.axvspan(i - 0.45, i + 0.45, alpha=0.12, color="deepskyblue", zorder=0)
ax.set_xticks(x_pos)
ax.set_xticklabels(disp, fontsize=TS - 2, rotation=15, ha="right")
ax.set_ylabel("dCl/dAoA_mod [/deg]", fontsize=TS)
ax.set_ylim(0.09, 0.18)
ax.set_title("線形域 Cl スロープの時系列変化（AoA_mod = −5〜+5°、線形回帰）", fontsize=TS)
ax.yaxis.set_major_locator(ticker.MultipleLocator(0.01))
ax.legend(fontsize=10)
ax.grid(axis="y", linestyle="--", alpha=0.5)
fig.tight_layout()
img_slope = fig_to_bytes(fig)

# ─── 図3: Cl 曲線オーバーレイ ─────────────────────────────────────────────────
fig, ax = plt.subplots(figsize=(7.5, 6))
for k, info in DATASETS.items():
    df = data[k]
    lw = 2.5 if "260605" in k else LW
    ms = 5   if "260605" in k else MK
    ax.plot(df["AoA_mod"], df["Cl"],
            color=info["color"], marker=info["marker"],
            linewidth=lw, markersize=ms,
            label=DISP_NAMES.get(k, k.replace("\n", " ")))
ax.axhline(0, color="black", linewidth=0.6, linestyle="--")
ax.axvline(0, color="black", linewidth=0.6, linestyle="--")
aoa_th = np.linspace(-25, 25, 200)
ax.plot(aoa_th, 2*np.pi*np.deg2rad(aoa_th), color="gray",
        linestyle=":", linewidth=1.5, label=r"$2\pi$ 理論")
ax.set_xlim(-25, 25)
ax.set_ylim(-1.0, 1.0)
ax.set_xlabel("AoA (補正後) [deg]", fontsize=TS)
ax.set_ylabel(r"$C_l$", fontsize=TS)
ax.set_title("Cl 曲線 全データ比較", fontsize=TS)
ax.legend(fontsize=8, ncol=2, loc="upper left")
ax.grid(linestyle="--", alpha=0.4)
fig.tight_layout()
img_cl = fig_to_bytes(fig)

# ─── 図4: AoA=0 付近の拡大 ───────────────────────────────────────────────────
fig, ax = plt.subplots(figsize=(7.5, 6))
for k, info in DATASETS.items():
    df = data[k]
    near = df[(df["AoA"] >= -5) & (df["AoA"] <= 5)].sort_values("AoA")
    lw = 2.5 if "260605" in k else LW
    ms = 6   if "260605" in k else 5
    ax.plot(near["AoA_mod"], near["Cl"],
            color=info["color"], marker=info["marker"],
            linewidth=lw, markersize=ms,
            label=DISP_NAMES.get(k, k.replace("\n", " ")))
ax.axhline(0, color="black", linewidth=1, linestyle="--")
ax.axvline(0, color="black", linewidth=0.6, linestyle="--")
ax.set_xlim(-5, 5)
ax.set_ylim(-0.5, 0.5)
ax.set_xlabel("AoA (補正後) [deg]", fontsize=TS)
ax.set_ylabel(r"$C_l$", fontsize=TS)
ax.set_title("AoA = −5〜+5° 拡大（零揚力迎角のずれ確認）", fontsize=TS)
ax.legend(fontsize=8, ncol=2)
ax.grid(linestyle="--", alpha=0.4)
fig.tight_layout()
img_zoom = fig_to_bytes(fig)

# ─── 図5: Cd・Cm 曲線オーバーレイ ────────────────────────────────────────────
fig, axes = plt.subplots(1, 2, figsize=(13, 5.5))
for k, info in DATASETS.items():
    df = data[k]
    lw = 2.5 if "260605" in k else LW
    ms = 5   if "260605" in k else MK
    lbl = DISP_NAMES.get(k, k.replace("\n", " "))
    axes[0].plot(df["AoA_mod"], df["Cd"],
                 color=info["color"], marker=info["marker"],
                 linewidth=lw, markersize=ms, label=lbl)
    axes[1].plot(df["AoA_mod"], df["Cm"],
                 color=info["color"], marker=info["marker"],
                 linewidth=lw, markersize=ms, label=lbl)
for ax, ylabel, ylim in zip(axes, [r"$C_d$", r"$C_m$"],
                                    [(0, 0.4), (-0.2, 0.2)]):
    ax.axhline(0, color="black", linewidth=0.6, linestyle="--")
    ax.axvline(0, color="black", linewidth=0.6, linestyle="--")
    ax.set_xlim(-25, 25)
    ax.set_ylim(*ylim)
    ax.set_xlabel("AoA (補正後) [deg]", fontsize=TS)
    ax.set_ylabel(ylabel, fontsize=TS)
    ax.legend(fontsize=7, ncol=2)
    ax.grid(linestyle="--", alpha=0.4)
axes[0].set_title("Cd 曲線 全データ比較", fontsize=TS)
axes[1].set_title("Cm 曲線 全データ比較", fontsize=TS)
fig.tight_layout()
img_cdcm = fig_to_bytes(fig)

# ─── 図6a/b: P-M対称性誤差 ───────────────────────────────────────────────────
pm_keys_raw  = [k for k in labels if not np.isnan(stats[k]["pm_rmse"])]
pm_keys_corr = [k for k in labels if not np.isnan(pm_rmse_corr[k])]

def bar_pm(keys, vals, title):
    fig, ax = plt.subplots(figsize=(max(5, len(keys)*0.9), 4.0))
    bc = [DATASETS[k]["color"] for k in keys]
    dp = [DISP_NAMES.get(k, k.replace("\n", " ")) for k in keys]
    bars = ax.bar(range(len(keys)), vals, color=bc, edgecolor="white", linewidth=0.8, width=0.6)
    for i, k in enumerate(keys):
        if "260605" in k:
            ax.axvspan(i - 0.4, i + 0.4, alpha=0.15, color="deepskyblue", zorder=0)
    for bar, v in zip(bars, vals):
        ax.text(bar.get_x() + bar.get_width()/2, v + 0.001,
                f"{v:.3f}", ha="center", va="bottom", fontsize=9, fontweight="bold")
    ax.set_xticks(range(len(keys)))
    ax.set_xticklabels(dp, fontsize=TS - 3, rotation=15, ha="right")
    ax.set_ylabel("Cl P-M RMSE", fontsize=TS)
    ax.set_title(title, fontsize=TS)
    ax.grid(axis="y", linestyle="--", alpha=0.5)
    fig.tight_layout()
    return fig_to_bytes(fig)

img_pm      = bar_pm(pm_keys_raw,
                     [stats[k]["pm_rmse"] for k in pm_keys_raw],
                     "（a）幾何迎角基準")
img_pm_corr = bar_pm(pm_keys_corr,
                     [pm_rmse_corr[k] for k in pm_keys_corr],
                     "（b）零揚力迎角基準（補正後）")

# ─── 図7: 260605 ② と ① の差 (Cl) ───────────────────────────────────────────
fig, axes = plt.subplots(1, 2, figsize=(12, 5))
# 左: 両データ重ね
k1, k2 = "260605①", "260605②"
for k, info in [(k1, DATASETS[k1]), (k2, DATASETS[k2])]:
    df = data[k]
    axes[0].plot(df["AoA_mod"], df["Cl"],
                 color=info["color"], marker=info["marker"],
                 linewidth=2, markersize=5,
                 label=DISP_NAMES[k])
axes[0].axhline(0, color="black", linewidth=0.6, linestyle="--")
axes[0].set_xlim(-30, 30); axes[0].set_ylim(-1.0, 1.0)
axes[0].set_xlabel("AoA (補正後) [deg]", fontsize=TS)
axes[0].set_ylabel(r"$C_l$", fontsize=TS)
axes[0].set_title("260605 ① vs ② 重ね合わせ", fontsize=TS)
axes[0].legend(fontsize=10); axes[0].grid(linestyle="--", alpha=0.4)

# 右: 差分
AOA_GRID = np.arange(-30, 31, dtype=float)
f1 = interp1d(data[k1]["AoA"].astype(float), data[k1]["Cl"].astype(float),
              kind="linear", bounds_error=False, fill_value=np.nan)
f2 = interp1d(data[k2]["AoA"].astype(float), data[k2]["Cl"].astype(float),
              kind="linear", bounds_error=False, fill_value=np.nan)
diff12 = f2(AOA_GRID) - f1(AOA_GRID)
axes[1].plot(AOA_GRID, diff12, color="steelblue", linewidth=2, marker="o", markersize=3)
axes[1].axhline(0, color="black", linewidth=0.8)
axes[1].set_xlim(-30, 30); axes[1].set_ylim(-0.05, 0.05)
axes[1].set_xlabel("AoA [deg]", fontsize=TS)
axes[1].set_ylabel(r"$\Delta C_l$ (② − ①)", fontsize=TS)
axes[1].set_title("260605 ② − ① の差分", fontsize=TS)
axes[1].yaxis.set_major_locator(ticker.MultipleLocator(0.01))
axes[1].grid(linestyle="--", alpha=0.4)
fig.tight_layout()
img_new_compare = fig_to_bytes(fig)

# ─── PowerPoint 作成 ──────────────────────────────────────────────────────────
prs = Presentation(TEMPLATE_PATH)
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

sldIdLst = prs.slides._sldIdLst
for sId in list(sldIdLst):
    rId = sId.get(_qn('r:id'))
    prs.part.drop_rel(rId)
    sldIdLst.remove(sId)

LAY_TITLE   = prs.slide_layouts[0]
LAY_CONTENT = prs.slide_layouts[1]

CA_L = Inches(0.333)
CA_T = Inches(1.167)
CA_W = Inches(9.333)
CA_H = Inches(5.535)

DARK    = RGBColor(0x1A, 0x1A, 0x2E)
ACC     = RGBColor(0x4A, 0x9E, 0xEA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
TBLHDR  = RGBColor(0x1A, 0x3A, 0x6A)
NOTE_BG = RGBColor(0xF0, 0xF4, 0xFF)

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=14, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def add_img(slide, img_bytes, l, t, w, h):
    img_bytes.seek(0)
    slide.shapes.add_picture(img_bytes, l, t, w, h)

def note_box(slide, text, l, t, w, h, size=12):
    add_rect(slide, l, t, w, h, fill=NOTE_BG, line=ACC, line_w=Pt(1))
    add_text(slide, text, l + Inches(0.1), t + Inches(0.1),
             w - Inches(0.2), h - Inches(0.2), size=size, color=DARK)

# ══════════════════════════════════════════════════════════════════
# スライド 1: タイトル
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_TITLE)
sl.placeholders[0].text = "NACA0012 Rigid Wing\n全データ比較（〜260610）／揚力傾斜の謎の解決"
sl.placeholders[1].text = "M1　岡本雄哉"
add_text(sl, ("Laboratory Seminar, Imamura Lab & Yamashita Lab\n"
              "The University of Tokyo, Engineering building 7,  June 5, 2026"),
         Inches(0.38), Inches(0.29), Inches(7.19), Inches(0.77),
         size=15, color=RGBColor(0x44, 0x44, 0x44))

# ══════════════════════════════════════════════════════════════════
# スライド 2: データセット一覧表
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_CONTENT)
sl.shapes.title.text = "データセット一覧（rigid 実験・自動更新）"

cols_w = [Inches(1.35), Inches(1.0), Inches(1.1), Inches(1.15),
          Inches(1.35), Inches(1.1), Inches(1.95)]
headers = ["データ", "日付", "Cl @ AoA=0", "零揚力迎角\n[deg]",
           "Clスロープ\n[/deg]", "P-M RMSE\n(Cl)", "備考"]
# 行数に応じて高さ・フォントを自動調整（件数が増えてもスライドに収める）
_n_rows  = len(stats) + 1                       # ヘッダ + データ行
row_h    = min(Inches(0.46), Emu(int(Inches(6.05) / _n_rows)))
_fs_hdr  = 10 if row_h >= Inches(0.36) else 8
_fs_cell = 10 if row_h >= Inches(0.36) else 8
top0     = CA_T

for ci, (hd, cw) in enumerate(zip(headers, cols_w)):
    lft = CA_L + sum(cols_w[:ci])
    add_rect(sl, lft, top0, cw, row_h, fill=TBLHDR,
             line=RGBColor(0x88, 0x99, 0xBB), line_w=Pt(0.5))
    add_text(sl, hd, lft + Inches(0.04), top0 + Inches(0.03),
             cw - Inches(0.08), row_h - Inches(0.06),
             size=_fs_hdr, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

stat_keys = list(stats.keys())
for ri, ((disp_name, who), k) in enumerate(zip(ROW_INFO, stat_keys)):
    s  = stats[k]
    bg = BG_COLOR[who]
    pm_str = f"{s['pm_rmse']:.3f}" if not np.isnan(s["pm_rmse"]) else "—"
    note   = NOTE_STR[who]
    vals   = [disp_name, s["date"],
              f"{s['Cl0']:+.4f}", f"{s['aoa0']:.3f}°",
              f"{s['slope']:.4f}", pm_str, note]
    for ci, (val, cw) in enumerate(zip(vals, cols_w)):
        lft = CA_L + sum(cols_w[:ci])
        top = top0 + row_h * (ri + 1)
        add_rect(sl, lft, top, cw, row_h, fill=bg,
                 line=RGBColor(0xCC, 0xCC, 0xCC), line_w=Pt(0.5))
        fsz = (_fs_cell - 1) if ci == 6 else _fs_cell
        add_text(sl, val, lft + Inches(0.04), top + Inches(0.04),
                 cw - Inches(0.08), row_h - Inches(0.08),
                 size=fsz, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# スライド 3: 零揚力迎角の時系列
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_CONTENT)
sl.shapes.title.text = "零揚力迎角の時系列変化"
add_img(sl, img_aoa0, CA_L, CA_T, Inches(6.5), Inches(4.5))

a260605_1 = stats["260605①"]["aoa0"]
a260605_2 = stats["260605②"]["aoa0"]
note_box(sl, (f"【時系列】\n"
              f"・伊東さん(2024): 0.23〜0.29°\n"
              f"・岡本0924/1020: 0.35〜0.54°\n"
              f"・岡本260430: 0.68°\n"
              f"・金光260520: 0.77°\n"
              f"・260605①(新): {a260605_1:.3f}°\n"
              f"・260605②(新): {a260605_2:.3f}°\n\n"
              f"→ 水色ハイライトが新システム\n"
              f"→ 過去データとの比較で\n"
              f"　取り付け角変化を確認"),
         CA_L + Inches(6.65), CA_T, Inches(2.6), Inches(4.5))

# ══════════════════════════════════════════════════════════════════
# スライド 4: Clスロープの時系列
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_CONTENT)
sl.shapes.title.text = "Cl スロープ（dCl/dAoA_mod）の時系列変化"
add_img(sl, img_slope, CA_L, CA_T, Inches(6.5), Inches(4.5))

sl260605_1 = stats["260605①"]["slope"]
sl260605_2 = stats["260605②"]["slope"]
note_box(sl, (f"【時系列】\n"
              f"・伊東さん/岡本: ~0.110〜0.112\n"
              f"・金光260417/0424: ~0.120〜0.122\n"
              f"・金光260520: 0.111(通常範囲)\n"
              f"・260605①(新): {sl260605_1:.4f}\n"
              f"・260605②(新): {sl260605_2:.4f}\n\n"
              f"点線: 2π/rad = 0.1097/deg\n\n"
              f"→ 新システムのスロープが\n"
              f"　既存データと整合するか確認"),
         CA_L + Inches(6.65), CA_T, Inches(2.6), Inches(4.5))

# ══════════════════════════════════════════════════════════════════
# スライド 5: Cl 曲線オーバーレイ
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_CONTENT)
sl.shapes.title.text = "Cl 曲線 全データ比較（260605 太線）"
add_img(sl, img_cl, CA_L, CA_T, Inches(6.5), Inches(5.2))
note_box(sl, ("【ポイント】\n"
              "・260605①② (水色・太線) が\n"
              "　過去データと重なるか確認\n\n"
              "・線形域での傾き一致\n"
              "・AoA=0 近傍のシフトは\n"
              "　次スライドで拡大確認\n\n"
              "・260605の2回計測は\n"
              "　新システムの再現性確認"),
         CA_L + Inches(6.65), CA_T, Inches(2.6), Inches(4.5))

# ══════════════════════════════════════════════════════════════════
# スライド 6: AoA=0 付近拡大
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_CONTENT)
sl.shapes.title.text = "AoA = −5〜+5° 拡大（0点付近・零揚力迎角確認）"
add_img(sl, img_zoom, CA_L, CA_T, Inches(6.5), Inches(5.2))
note_box(sl, ("【読み方】\n"
              "Cl=0 となる AoA が\n"
              "データによって異なる\n\n"
              "伊東さん: 0.23〜0.29°\n"
              "岡本:     0.35〜0.68°\n"
              "金光くん: 0.77〜0.81°\n\n"
              f"260605①: {a260605_1:.3f}°\n"
              f"260605②: {a260605_2:.3f}°\n\n"
              "→ 新システムの零揚力迎角が\n"
              "　直前の計測と比較して\n"
              "　変化しているか確認"),
         CA_L + Inches(6.65), CA_T, Inches(2.6), Inches(5.0))

# ══════════════════════════════════════════════════════════════════
# スライド 7: Cd・Cm 曲線
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_CONTENT)
sl.shapes.title.text = "Cd・Cm 曲線 全データ比較（260605 太線）"
add_img(sl, img_cdcm, CA_L, CA_T, CA_W, Inches(5.0))

# ══════════════════════════════════════════════════════════════════
# スライド 8: P-M 対称性誤差
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_CONTENT)
sl.shapes.title.text = "正負掃引の非対称性（P-M 対称性誤差）"
GW, GH = Inches(4.55), Inches(3.5)
add_img(sl, img_pm,      CA_L,                CA_T, GW, GH)
add_img(sl, img_pm_corr, CA_L + Inches(4.75), CA_T, GW, GH)
note_box(sl, ("→ 零揚力迎角を基準に補正すると（b）、取り付け角オフセット α₀ による\n"
              "　見かけ上の非対称性が除去され、実質的な空力非対称性を評価できる。\n"
              "→ 260605 の P-M RMSE が過去データと同程度なら新システムの信頼性を確認。"),
         CA_L, CA_T + Inches(3.65), CA_W, Inches(1.75), size=12)

# ══════════════════════════════════════════════════════════════════
# スライド 9: 260605 ① vs ② 比較（新システム再現性）
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_CONTENT)
sl.shapes.title.text = "260605 ①② 比較：新システムの再現性確認"
add_img(sl, img_new_compare, CA_L, CA_T, CA_W, Inches(5.0))
note_box(sl, ("→ 同日に2回計測した ① と ② の Cl 差分。\n"
              "　|ΔCl| が小さければ新システムの計測再現性が高いと評価できる。"),
         CA_L, CA_T + Inches(5.1), CA_W, Inches(0.9), size=12)

# ══════════════════════════════════════════════════════════════════
# スライド 10: まとめ
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_CONTENT)
sl.shapes.title.text = "まとめ"

sl_r3 = stats["rigid3\n修正後"]["slope"]
a_r3  = stats["rigid3\n修正後"]["aoa0"]
items = [
    ("① 揚力傾斜の異常 = 差圧デジボルの接触不良が真因",
     f"260605①②/260608 の Cl スロープが過去(0.110〜0.112/deg)より大幅に増（260605①={sl260605_1:.3f}/deg 等）。\n"
     "原因：差圧電圧(mV)が接触不良で過小 → 風速U・動圧qが過小 → Cl=L/q が過大 → 傾斜が過大に見えていた。"),
    ("② 接触修正後 rigid3 で旧データと一致（解決確認）",
     f"接触を修正した rigid3（真mV≈1320, U≈13.1m/s）の傾斜 = {sl_r3:.4f}/deg ≒ 過去の通常範囲(0.110〜0.112)・2π。\n"
     "→ 同じ翼なら傾斜は風速によらず一定、という物理がようやく再現された。揚力傾斜の謎は解決。"),
    ("③ 零揚力迎角 α₀ は別問題（取り付け角）",
     f"α₀は 260605({a260605_1:.2f}/{a260605_2:.2f}°)→rigid3({a_r3:.2f}°) と再マウントで旧(0.5〜0.8°)へ収束。\n"
     "α₀は q補正の影響を受けない（Cl=0の位置は不変）ため、傾斜(デジボル)とは独立した別原因だった。"),
    ("④ 今後の方針",
     "計測前チェックに『差圧mVの妥当性確認（デジボル接触）』を追加。過去の新データは真mVが分かれば補正可能。\n"
     "疑っていたマウント弾性/センサFy軸/定格/後処理は全て無罪と確認済み。"),
]

ITEM_H   = Inches(1.2)
ITEM_GAP = Inches(0.06)
for i, (title, body) in enumerate(items):
    top = CA_T + (ITEM_H + ITEM_GAP) * i
    add_rect(sl, CA_L, top, CA_W, ITEM_H,
             fill=RGBColor(0xE8, 0xF0, 0xFE), line=ACC, line_w=Pt(1))
    add_rect(sl, CA_L, top, Inches(0.07), ITEM_H, fill=ACC)
    add_text(sl, title, CA_L + Inches(0.15), top + Inches(0.04),
             CA_W - Inches(0.2), Inches(0.32), size=13, bold=True, color=DARK)
    add_text(sl, body,  CA_L + Inches(0.15), top + Inches(0.38),
             CA_W - Inches(0.2), Inches(0.76), size=10, color=DARK)

# ─── 保存 ─────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"保存完了: {OUT}")
