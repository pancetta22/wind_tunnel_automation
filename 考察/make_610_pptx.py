#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""260610 実験結果まとめ（簡潔版）スライド生成。揚力傾斜は ±5° で算出。"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

C_DARK=RGBColor(0x1A,0x37,0x6C); C_MID=RGBColor(0x2E,0x6E,0xB4)
C_GREEN=RGBColor(0x21,0x7A,0x3C); C_GREEN_L=RGBColor(0xD6,0xEE,0xDF)
C_ORANGE=RGBColor(0xD9,0x6D,0x00); C_ORANGE_L=RGBColor(0xFD,0xF0,0xDE)
C_GRAY=RGBColor(0xF4,0xF6,0xF9); C_TEXT=RGBColor(0x1A,0x1A,0x2E)
C_WHITE=RGBColor(0xFF,0xFF,0xFF); C_RED=RGBColor(0xC0,0x39,0x2B); C_SUB=RGBColor(0x5A,0x6A,0x85)

prs=Presentation(); prs.slide_width=Inches(13.33); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]

def rect(s,l,t,w,h,fill=None,border=None,bw=Pt(1)):
    sh=s.shapes.add_shape(1,l,t,w,h); sh.shadow.inherit=False
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    else: sh.fill.background()
    if border: sh.line.color.rgb=border; sh.line.width=bw
    else: sh.line.fill.background()
    return sh

def text(s,t,l,tp,w,h,size=Pt(14),bold=False,color=C_TEXT,align=PP_ALIGN.LEFT,it=False):
    tb=s.shapes.add_textbox(l,tp,w,h); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=t
    r.font.size=size; r.font.bold=bold; r.font.color.rgb=color; r.font.italic=it
    return tb

def header(s,title,sub=None):
    rect(s,0,Inches(0.85),prs.slide_width,prs.slide_height-Inches(0.85),fill=C_GRAY)
    rect(s,0,0,prs.slide_width,Inches(0.85),fill=C_DARK)
    text(s,title,Inches(0.35),Inches(0.08),Inches(12.5),Inches(0.55),size=Pt(23),bold=True,color=C_WHITE)
    if sub: text(s,sub,Inches(0.37),Inches(0.56),Inches(12),Inches(0.3),size=Pt(11.5),color=RGBColor(0xA8,0xC8,0xF0))

def bullets(s,lines,l,t,w,h,size=Pt(13),gap=Pt(6)):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True; first=True
    for it in lines:
        txt=it[0]; lv=it[1] if len(it)>1 else 0; col=it[2] if len(it)>2 else C_TEXT; bd=it[3] if len(it)>3 else False
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.space_before=gap; p.level=lv
        mark="■ " if lv==0 else "・ "
        r=p.add_run(); r.text=mark+txt; r.font.size=size if lv==0 else Pt(size.pt-1.5)
        r.font.color.rgb=col; r.font.bold=bd

def panel(s,l,t,w,h,fill,border): rect(s,l,t,w,h,fill=fill,border=border,bw=Pt(1.5))

# ===== S1 サマリ =====
s=prs.slides.add_slide(BLANK)
header(s,"260610 実験結果：新セットアップが安定（0608を再現）","揚力傾斜は ±5° で算出。再マウント後の状態が再現性を持って定着")

for (l,title,val,unit,sub) in [
 (0.4,"揚力傾斜 dCl/dα","4.49","/rad","旧3.94 / 605=4.92 / 608=4.44"),
 (4.6,"ゼロ揚力角 α₀","0.75","°","旧0.5〜0.7と一致（605=1.44）"),
 (8.8,"最大揚力 / 失速角","0.947","@15°","失速角が旧と同じ15°に復帰")]:
    panel(s,Inches(l),Inches(1.1),Inches(4.0 if l<8 else 4.1),Inches(2.5),C_GREEN_L,C_GREEN)
    text(s,title,Inches(l+0.15),Inches(1.2),Inches(3.7),Inches(0.4),size=Pt(13),bold=True,color=C_GREEN)
    text(s,val,Inches(l+0.15),Inches(1.65),Inches(3.7),Inches(0.9),size=Pt(40),bold=True,color=C_DARK)
    text(s,unit,Inches(l+0.15),Inches(2.62),Inches(3.7),Inches(0.4),size=Pt(13),color=C_SUB)
    text(s,sub,Inches(l+0.15),Inches(2.98),Inches(3.8),Inches(0.5),size=Pt(11),color=C_SUB)

panel(s,Inches(0.4),Inches(3.8),Inches(12.5),Inches(3.05),C_WHITE,C_MID)
text(s,"ポイント",Inches(0.55),Inches(3.9),Inches(12),Inches(0.4),size=Pt(15),bold=True,color=C_DARK)
bullets(s,[
 ("0610は0608をほぼ再現（傾斜 4.49≈4.44）→ 再マウント後の状態が安定・再現性あり",0,C_GREEN,True),
 ("α₀=0.75°（旧と一致）、失速角15°（旧と一致）、Cd0=0.0155（旧へ接近）→ 多くの指標が旧と整合",0),
 ("605→608の大改善（1.25→1.13×）は再マウント。608→610はほぼ同じ＝定着",0,C_SUB),
 ("ただし揚力傾斜の残差 1.14× は安定して残る（消えない系統残差）",0,C_ORANGE,True),
],Inches(0.55),Inches(4.35),Inches(12.1),Inches(2.4),size=Pt(13.5),gap=Pt(9))

# ===== S2 推移表 =====
s=prs.slides.add_slide(BLANK)
header(s,"605 → 608 → 610 の推移（剛体 NACA0012, ±5°）")
data=[
 ("旧0520","1165","12.23","3.93","0.72","0.808@15","—",False),
 ("旧0430","1166","12.08","3.91","0.68","0.812@15","—",False),
 ("旧1020","1163","12.11","3.97","0.52","0.791@15","—",False),
 ("新0605","1056","11.65","4.92","1.44","1.045@16","1.25×",False),
 ("新0608","1171","12.29","4.44","0.93","0.914@16","1.13×",False),
 ("新0610 ★","1161","12.20","4.49","0.75","0.947@15","1.14×",True),
]
heads=["実験","mV","U[m/s]","dCl/dα[/rad]","α₀[°]","Cl_max@°","旧比"]
xs=[0.5,2.6,3.7,5.0,7.3,8.6,10.2]; ws=[2.1,1.1,1.3,2.3,1.3,1.6,1.6]
y=Inches(1.2)
rect(s,Inches(0.4),y,Inches(12.5),Inches(0.42),fill=C_DARK)
for hx,hw,ht in zip(xs,ws,heads):
    text(s,ht,Inches(hx),y+Inches(0.04),Inches(hw),Inches(0.34),size=Pt(12.5),bold=True,color=C_WHITE)
y=Inches(1.66)
for i,(lab,mv,U,rad,a0,clm,ratio,hl) in enumerate(data):
    bg=C_ORANGE_L if hl else (C_WHITE if i%2==0 else RGBColor(0xEC,0xF1,0xF8))
    rect(s,Inches(0.4),y,Inches(12.5),Inches(0.62),fill=bg,border=(C_ORANGE if hl else RGBColor(0xD0,0xDA,0xE8)),bw=Pt(2 if hl else 0.75))
    col=C_RED if hl else C_TEXT
    for hx,hw,v in zip(xs,ws,[lab,mv,U,rad,a0,clm,ratio]):
        text(s,v,Inches(hx),y+Inches(0.13),Inches(hw),Inches(0.4),size=Pt(13),bold=hl,color=col)
    y+=Inches(0.62)
panel(s,Inches(0.4),Inches(5.7),Inches(12.5),Inches(1.15),C_GREEN_L,C_GREEN)
text(s,"傾斜は 605(4.92) → 608(4.44) → 610(4.49) で 608/610 に定着。α₀ は 1.44 → 0.93 → 0.75 と旧へ収束、失速角も15°復帰。",
     Inches(0.6),Inches(5.95),Inches(12.1),Inches(0.6),size=Pt(13),bold=True,color=C_DARK)

# ===== S3 考察 =====
s=prs.slides.add_slide(BLANK)
header(s,"考察：残差 1.14× の正体")
panel(s,Inches(0.4),Inches(1.1),Inches(6.15),Inches(3.5),C_WHITE,C_MID)
text(s,"α₀・失速・Cd0 は一致、傾斜だけ高い",Inches(0.55),Inches(1.2),Inches(5.85),Inches(0.4),size=Pt(15),bold=True,color=C_DARK)
bullets(s,[
 ("α₀=0.75°（旧と一致）：取付角は正しく合った",0,C_GREEN),
 ("失速角15°・Cd0も旧へ接近：粗いズレは解消",0,C_GREEN),
 ("なのに揚力傾斜だけ 14% 高い",0,C_RED,True),
 ("→ 「取付角」では説明できない残差",0),
],Inches(0.55),Inches(1.65),Inches(5.85),Inches(2.8),size=Pt(13))

panel(s,Inches(6.75),Inches(1.1),Inches(6.15),Inches(3.5),C_ORANGE_L,C_ORANGE)
text(s,"最有力：マウントのねじれ柔性（空力弾性）",Inches(6.9),Inches(1.2),Inches(5.85),Inches(0.4),size=Pt(14),bold=True,color=C_ORANGE)
bullets(s,[
 ("支持系が荷重でねじれ、実効迎角が増える",0,C_RED,True),
 ("揚力ゼロ(α₀)では荷重ゼロ→ねじれゼロ→α₀は一致",0),
 ("迎角UP→荷重UP→ねじれUP→傾斜だけ急になる",0),
 ("手掛け水平荷重テストの「撓み・ヒステリシス」とも整合",0),
 ("再マウントで剛性変化→605→608で傾斜が動いた事実と一致",0),
],Inches(6.9),Inches(1.65),Inches(5.85),Inches(2.8),size=Pt(12.5),gap=Pt(6))

panel(s,Inches(0.4),Inches(4.8),Inches(12.5),Inches(2.0),RGBColor(0x14,0x2A,0x55),C_MID)
text(s,"決定的な次のテスト",Inches(0.6),Inches(4.92),Inches(11),Inches(0.4),size=Pt(15),bold=True,color=C_WHITE)
bullets(s,[
 ("同一マウントのまま 風速2点 で揚力傾斜を比較",0,C_WHITE,True),
 ("空力弾性なら 動圧qが高いほど傾斜が急（発散的）→ 確定",0,RGBColor(0xC8,0xF0,0xD8)),
 ("確定後：支持系の剛性向上 or 補正、または『どちらが真値か』の再評価（2D理論は2π≈6.28/rad）",0,RGBColor(0xC8,0xF0,0xD8)),
],Inches(0.6),Inches(5.35),Inches(12.1),Inches(1.4),size=Pt(13),gap=Pt(8))

OUT=os.path.join(os.path.dirname(os.path.abspath(__file__)),"Windy_260610結果.pptx")
prs.save(OUT)
print("生成完了:",OUT,"(%dスライド)"%len(prs.slides._sldIdLst))
