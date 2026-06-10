#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Windy フォルダ構成ガイド（操作マニュアル）PowerPoint 生成スクリプト

リポジトリ全体のフォルダ・ファイル構成と役割、データの流れを説明する。
レイアウト崩れを防ぐため、全テキスト枠に「はみ出し時の自動縮小」を適用。
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# ============================================================
#  カラーパレット
# ============================================================
C_DARK_BLUE   = RGBColor(0x1A, 0x37, 0x6C)   # タイトルバー背景
C_MID_BLUE    = RGBColor(0x2E, 0x6E, 0xB4)   # アクセント
C_GREEN       = RGBColor(0x21, 0x7A, 0x3C)
C_GREEN_LIGHT = RGBColor(0xD6, 0xEE, 0xDF)
C_ORANGE      = RGBColor(0xD9, 0x6D, 0x00)
C_ORANGE_LIGHT= RGBColor(0xFD, 0xF0, 0xDE)
C_GRAY_BG     = RGBColor(0xF4, 0xF6, 0xF9)
C_PROG_BG     = RGBColor(0xE8, 0xF0, 0xFB)
C_TEXT        = RGBColor(0x1A, 0x1A, 0x2E)
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_ROW_ALT     = RGBColor(0xEE, 0xF3, 0xFA)   # ゼブラ行
C_ROW_BORDER  = RGBColor(0xD2, 0xDC, 0xEA)
C_PURPLE      = RGBColor(0x6A, 0x3D, 0x9A)
C_TEAL        = RGBColor(0x10, 0x80, 0x80)
C_GRAYTXT     = RGBColor(0x66, 0x66, 0x66)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ============================================================
#  ヘルパー
# ============================================================
def _shrink(tf):
    """テキストが枠を超えたら自動縮小（レイアウト崩れ防止）。"""
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    if bodyPr is None:
        return
    for tag in ('a:normAutofit', 'a:spAutoFit', 'a:noAutofit'):
        e = bodyPr.find(qn(tag))
        if e is not None:
            bodyPr.remove(e)
    bodyPr.append(bodyPr.makeelement(qn('a:normAutofit'), {}))


def add_rect(slide, l, t, w, h, fill=None, border=None, border_w=Pt(1)):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = border_w
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=Pt(14), bold=False, color=C_TEXT,
             align=PP_ALIGN.LEFT, italic=False, wrap=True, anchor=None):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(1);  tf.margin_bottom = Pt(1)
    if anchor is not None:
        tf.vertical_anchor = anchor
    lines = text.split('\n')
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = ln
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.italic = italic
    _shrink(tf)
    return tb


def slide_header(slide, title, subtitle=None):
    add_rect(slide, 0, Inches(0.85), SW, SH - Inches(0.85), fill=C_GRAY_BG)
    add_rect(slide, 0, 0, SW, Inches(0.85), fill=C_DARK_BLUE)
    add_rect(slide, 0, Inches(0.85), SW, Inches(0.06), fill=C_MID_BLUE)
    add_text(slide, title, Inches(0.35), Inches(0.08), Inches(10.5), Inches(0.5),
             font_size=Pt(23), bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.37), Inches(0.55), Inches(11), Inches(0.3),
                 font_size=Pt(12), color=RGBColor(0xA8, 0xC8, 0xF0))


def tag_box(slide, text, l, t, color, w=Inches(2.3)):
    add_rect(slide, l, t, w, Inches(0.34), fill=color)
    add_text(slide, text, l, t + Inches(0.02), w, Inches(0.30),
             font_size=Pt(11), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


def file_table(slide, rows, l, t, w, name_w, name_color,
               row_h=Inches(0.6), name_size=Pt(13), desc_size=Pt(12)):
    """ファイル一覧（filename | 役割）をゼブラ模様の表で描画。"""
    from pptx.enum.text import MSO_ANCHOR
    for i, (name, desc) in enumerate(rows):
        ry = t + i * row_h
        bg = C_WHITE if i % 2 == 0 else C_ROW_ALT
        add_rect(slide, l, ry, w, row_h - Inches(0.04),
                 fill=bg, border=C_ROW_BORDER, border_w=Pt(0.75))
        add_text(slide, name, l + Inches(0.12), ry, name_w, row_h - Inches(0.04),
                 font_size=name_size, bold=True, color=name_color,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, desc, l + name_w + Inches(0.2), ry,
                 w - name_w - Inches(0.35), row_h - Inches(0.04),
                 font_size=desc_size, color=C_TEXT, anchor=MSO_ANCHOR.MIDDLE)


def note_bar(slide, label, lines, l, t, w, h, color, bg):
    add_rect(slide, l, t, w, h, fill=bg, border=color, border_w=Pt(1.3))
    add_rect(slide, l, t, Inches(0.08), h, fill=color)
    add_text(slide, label, l + Inches(0.18), t + Inches(0.06), w - Inches(0.3), Inches(0.3),
             font_size=Pt(12), bold=True, color=color)
    add_text(slide, '\n'.join(lines), l + Inches(0.18), t + Inches(0.36),
             w - Inches(0.36), h - Inches(0.44),
             font_size=Pt(11.5), color=C_TEXT)


# ============================================================
#  スライド 1: タイトル
# ============================================================
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SW, SH, fill=C_GRAY_BG)
add_rect(slide, 0, 0, SW, Inches(4.7), fill=C_DARK_BLUE)
add_rect(slide, 0, Inches(4.7), SW, Inches(0.08), fill=C_MID_BLUE)

add_text(slide, "Windy 風洞実験自動計測システム",
         Inches(1.0), Inches(1.25), Inches(11.3), Inches(1.0),
         font_size=Pt(40), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "フォルダ構成ガイド",
         Inches(1.0), Inches(2.55), Inches(11.3), Inches(0.8),
         font_size=Pt(26), color=RGBColor(0xA8, 0xC8, 0xF0), align=PP_ALIGN.CENTER)
add_rect(slide, Inches(3.0), Inches(3.75), Inches(7.33), Inches(0.55), fill=C_MID_BLUE)
add_text(slide, "どのフォルダに何があるか・データはどう流れるか",
         Inches(3.0), Inches(3.82), Inches(7.33), Inches(0.42),
         font_size=Pt(15), color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "MATLAB + Python  /  迎角ステージ・6軸センサ・差圧デジボルを統合自動化",
         Inches(1.0), Inches(5.25), Inches(11.3), Inches(0.6),
         font_size=Pt(16), color=C_TEXT, align=PP_ALIGN.CENTER, italic=True)


# ============================================================
#  スライド 2: 全体像（リポジトリツリー）
# ============================================================
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "全体像：リポジトリ構成", "ルート直下は「入口」だけ。中身は役割ごとにフォルダ分け")

# 左：ツリー
add_rect(slide, Inches(0.4), Inches(1.15), Inches(7.0), Inches(5.95),
         fill=C_WHITE, border=C_ROW_BORDER, border_w=Pt(1))
tree = [
    ("Windy/",                       C_DARK_BLUE, True,  Pt(14)),
    ("├ run_experiment.m   ← 実験はこれを実行", C_GREEN, True, Pt(12.5)),
    ("├ run_postprocess.m  ← 後処理だけ再実行", C_GREEN, True, Pt(12.5)),
    ("├ setup_paths.m      ← 診断ツール用パス", C_GREEN, True, Pt(12.5)),
    ("├ config.json(.example)  設定ファイル",   C_TEXT,  False, Pt(12.5)),
    ("├ README.md / SPEC.md",                  C_TEXT,  False, Pt(12.5)),
    ("│",                            C_GRAYTXT,  False, Pt(12.5)),
    ("├ measurement_control/  計測機器の制御", C_MID_BLUE, True, Pt(12.5)),
    ("├ diagnostics/          点検・診断ツール", C_MID_BLUE, True, Pt(12.5)),
    ("├ leptrino/             6軸センサ通信",   C_MID_BLUE, True, Pt(12.5)),
    ("├ post_process/         後処理（Python）", C_MID_BLUE, True, Pt(12.5)),
    ("├ analysis/             結果の比較・分析", C_MID_BLUE, True, Pt(12.5)),
    ("└ manual/               このマニュアル",  C_MID_BLUE, True, Pt(12.5)),
]
for i, (txt, col, bd, sz) in enumerate(tree):
    add_text(slide, txt, Inches(0.65), Inches(1.3) + i * Inches(0.42),
             Inches(6.6), Inches(0.4), font_size=sz, bold=bd, color=col)

# 右：3グループ要約
gx = Inches(7.7)
groups = [
    ("🟢 ルート直下＝入口", C_GREEN, C_GREEN_LIGHT,
     ["実験・後処理・設定。普段触るのはここ。",
      "run_experiment / run_postprocess / config.json"]),
    ("🔵 計測ロジック", C_MID_BLUE, C_PROG_BG,
     ["機器制御・センサ通信・点検ツール。",
      "measurement_control / diagnostics / leptrino"]),
    ("🟣 後処理・分析", C_PURPLE, RGBColor(0xEE, 0xE7, 0xF6),
     ["空力係数の算出と過去データとの比較。",
      "post_process / analysis / manual"]),
]
gy = Inches(1.25)
for label, col, bg, lines in groups:
    note_bar(slide, label, lines, gx, gy, Inches(5.2), Inches(1.75), col, bg)
    gy += Inches(1.95)


# ============================================================
#  スライド 3: ルート直下のファイル
# ============================================================
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "ルート直下のファイル", "普段の操作で使う「入口」。MATLAB コマンドウィンドウから実行")

rows = [
    ("run_experiment.m",  "メイン実験スクリプト。計測〜後処理まで一気通貫（まずこれ）"),
    ("run_postprocess.m", "後処理だけを単体で再実行  例) run_postprocess('…/260605_rigid')"),
    ("setup_paths.m",     "診断ツールやヘルパを単体で使う前に1回だけ実行（パスを通す）"),
    ("config.json",       "各自の環境設定（COMポート・Pythonパス・保存先）※Git管理外"),
    ("config.json.example","設定の雛形。コピーして config.json を作る"),
    ("README.md",         "クイックスタート・構成・使い方の説明書"),
    ("SPEC.md",           "設計仕様書（実装の経緯・歴史的資料）"),
]
file_table(slide, rows, Inches(0.4), Inches(1.25), Inches(12.53), Inches(3.9),
           C_GREEN, row_h=Inches(0.66))

note_bar(slide, "💡 最短の使い方",
         ["① config.json.example をコピーして config.json を作り、COMポート等を設定（初回のみ）",
          "② MATLAB でこのフォルダに cd して  run_experiment  を実行 → あとは画面の指示どおり"],
         Inches(0.4), Inches(5.95), Inches(12.53), Inches(1.05), C_GREEN, C_GREEN_LIGHT)


# ============================================================
#  スライド 4: measurement_control/
# ============================================================
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "measurement_control/  — 計測機器の制御",
             "run_experiment が内部で使う MATLAB ヘルパ群")
tag_box(slide, "⚙ run_experiment が自動で使用", Inches(10.0), Inches(0.25), C_MID_BLUE, Inches(3.0))

rows = [
    ("QT_ADL1.m",        "迎角ステージ（QT-ADL1）のドライバクラス。原点復帰・角度移動"),
    ("LeptrinoLogger.m", "Leptrino 6軸センサの時系列ロガークラス（バックグラウンド記録）"),
    ("WindyMonitor.m",   "計測中のリアルタイムモニタ画面（波形・進捗・停止ボタン）"),
    ("get_sensor_data.m","6軸センサの瞬時平均値を1回取得する関数"),
    ("get_voltage.m",    "R6441B デジボル（差圧電圧）を取得する関数"),
    ("make_filename.m",  "計測ファイル名を規則に従って生成するユーティリティ"),
]
file_table(slide, rows, Inches(0.4), Inches(1.3), Inches(12.53), Inches(3.7),
           C_MID_BLUE, row_h=Inches(0.7))

note_bar(slide, "📌 ポイント",
         ["run_experiment 実行時に自動でパスが通るので、普段はここを直接触らない。",
          "個別に使いたいときは先に setup_paths を実行する。"],
         Inches(0.4), Inches(5.85), Inches(12.53), Inches(1.1), C_MID_BLUE, C_PROG_BG)


# ============================================================
#  スライド 5: diagnostics/
# ============================================================
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "diagnostics/  — 点検・診断ツール",
             "実験前後に手動で実行して機器やセンサを確認する")
tag_box(slide, "👤 手動で実行（要 setup_paths）", Inches(9.9), Inches(0.25), C_GREEN, Inches(3.1))

rows = [
    ("QT_ADL1_check_connection.m", "迎角ステージの COM ポート接続・通信を確認"),
    ("check_sensor_limit.m",       "6軸センサの定格（最大計測レンジ）を確認"),
    ("weight_check.m",             "既知のおもりを載せて力センサの読みを検証"),
    ("tare_measure.m",             "ゼロ点を取り、その基準からの6軸力を表示"),
    ("lumix_check_connection.py",  "カメラ（LUMIX）の接続確認（Python）"),
]
file_table(slide, rows, Inches(0.4), Inches(1.3), Inches(12.53), Inches(4.5),
           C_GREEN, row_h=Inches(0.72))

note_bar(slide, "👤 使い方",
         ["MATLAB のルートで  setup_paths  を1回実行 → 各ツールを関数名で実行できる。",
          "いずれも config.json と leptrino/ をリポジトリルート基準で参照する。"],
         Inches(0.4), Inches(5.95), Inches(12.53), Inches(1.05), C_GREEN, C_GREEN_LIGHT)


# ============================================================
#  スライド 6: leptrino/ と post_process/
# ============================================================
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "leptrino/  と  post_process/",
             "センサ通信（Python）と、計測後の数値処理（Python）")

# 左: leptrino
add_rect(slide, Inches(0.4), Inches(1.2), Inches(6.1), Inches(5.8),
         fill=C_WHITE, border=C_ROW_BORDER, border_w=Pt(1))
add_text(slide, "leptrino/", Inches(0.6), Inches(1.32), Inches(5.7), Inches(0.4),
         font_size=Pt(17), bold=True, color=C_TEAL)
add_text(slide, "6軸センサとの通信（32bit Python 経由）",
         Inches(0.6), Inches(1.74), Inches(5.7), Inches(0.32),
         font_size=Pt(12), color=C_GRAYTXT)
file_table(slide, [
    ("leptrino_server.py", "センサ計測サーバ。MATLABから呼ばれCSVへ記録"),
    ("CfsUsb.dll",         "Leptrino USB ドライバ（32bit専用）"),
], Inches(0.55), Inches(2.2), Inches(5.8), Inches(3.7), C_TEAL,
   row_h=Inches(1.0), name_size=Pt(12.5))
note_bar(slide, "⚠ 32bit 必須",
         ["CfsUsb.dll は 32bit 専用。config.json の",
          "python_exe には 32bit Python を指定する。"],
         Inches(0.55), Inches(4.4), Inches(5.8), Inches(1.2), C_ORANGE, C_ORANGE_LIGHT)

# 右: post_process
add_rect(slide, Inches(6.7), Inches(1.2), Inches(6.23), Inches(5.8),
         fill=C_WHITE, border=C_ROW_BORDER, border_w=Pt(1))
add_text(slide, "post_process/", Inches(6.9), Inches(1.32), Inches(5.8), Inches(0.4),
         font_size=Pt(17), bold=True, color=C_MID_BLUE)
add_text(slide, "計測データ → 風速・空力係数・グラフ（64bit Python）",
         Inches(6.9), Inches(1.74), Inches(5.9), Inches(0.32),
         font_size=Pt(12), color=C_GRAYTXT)
file_table(slide, [
    ("make_windspeed.py", "差圧電圧 → 風速 windspeed.csv"),
    ("calc_force.py",     "6軸力 → 空力係数 C_aero.csv・グラフPNG"),
    ("requirements.txt",  "必要 Python パッケージ一覧"),
    ("venv/",             "自動生成される仮想環境（Git管理外）"),
], Inches(6.85), Inches(2.2), Inches(5.95), Inches(3.7), C_MID_BLUE,
   row_h=Inches(0.62), name_size=Pt(12.5))
note_bar(slide, "⚙ 自動実行",
         ["run_experiment 完了時に venv 構築から自動。",
          "失敗時は run_postprocess('…') で再実行。"],
         Inches(6.85), Inches(4.85), Inches(5.95), Inches(1.15), C_MID_BLUE, C_PROG_BG)


# ============================================================
#  スライド 7: analysis/
# ============================================================
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "analysis/  — 結果の比較・分析",
             "新システムの結果を過去の剛体翼データと比較（パワポ自動更新）")
tag_box(slide, "🟣 rigid 実験で自動更新", Inches(10.3), Inches(0.25), C_PURPLE, Inches(2.7))

rows = [
    ("update_aero_data.py",            "新実験の C_aero.csv を取り込み→比較パワポ再生成（これ1つでOK）"),
    ("make_rigid_comparison_local.py", "比較パワポを生成する本体スクリプト"),
    ("Windy新システムによる実験結果.pptx", "比較パワポ（成果物）。実験追加で自動更新される"),
    ("研究室MTGテンプレート.pptx",      "パワポの雛形（研究室フォーマット）"),
    ("aero_data/",                     "各実験の空力係数データ C_aero.csv の置き場"),
    ("archive/",                       "使い終わった単発スクリプト・旧パワポ・グラフ"),
]
file_table(slide, rows, Inches(0.4), Inches(1.3), Inches(12.53), Inches(3.95),
           C_PURPLE, row_h=Inches(0.66))

note_bar(slide, "🟣 流れ",
         ["実験名に rigid を含むと、後処理の最後に「過去データと比較しますか？」→ y で自動更新。",
          "手動なら  python update_aero_data.py  （引数なしで config.json の output_dir を走査）。"],
         Inches(0.4), Inches(5.95), Inches(12.53), Inches(1.05), C_PURPLE, RGBColor(0xEE, 0xE7, 0xF6))


# ============================================================
#  スライド 8: データの流れ
# ============================================================
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "データの流れ", "計測から比較パワポまで、ファイルがどう生まれるか")

flow = [
    ("run_experiment.m", "MATLAB で実行\n4フェーズを自動計測", C_GREEN_LIGHT, C_GREEN),
    ("実験フォルダ", "output_dir/<実験名>/\ndata/・volt_summary・log", C_GRAY_BG, C_MID_BLUE),
    ("post_process", "windspeed.csv\nC_aero.csv・グラフPNG", C_PROG_BG, C_MID_BLUE),
    ("analysis", "比較パワポを更新\n（rigid 実験のみ）", RGBColor(0xEE, 0xE7, 0xF6), C_PURPLE),
]
bw, bh, gap = Inches(2.7), Inches(1.85), Inches(0.55)
total = len(flow) * bw + (len(flow) - 1) * gap
sx = (SW - total) / 2
y = Inches(1.8)
for i, (ttl, body, bg, fg) in enumerate(flow):
    x = sx + i * (bw + gap)
    add_rect(slide, x, y, bw, bh, fill=bg, border=fg, border_w=Pt(2.5))
    add_rect(slide, x, y, bw, Inches(0.55), fill=fg)
    add_text(slide, ttl, x, y + Inches(0.05), bw, Inches(0.45),
             font_size=Pt(15), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, body, x + Inches(0.1), y + Inches(0.68), bw - Inches(0.2), Inches(1.05),
             font_size=Pt(12.5), color=C_TEXT, align=PP_ALIGN.CENTER)
    if i < len(flow) - 1:
        add_text(slide, "▶", x + bw, y + Inches(0.62), gap, Inches(0.6),
                 font_size=Pt(24), bold=True, color=C_MID_BLUE, align=PP_ALIGN.CENTER)

note_bar(slide, "📁 実験フォルダの中身（output_dir/<実験名>/）",
         ["data/ … 各計測点の6軸センサCSV     ＊_volt_summary.csv … フェーズ毎の差圧電圧",
          "＊_experiment_log.json … 気温・気圧・校正定数      windspeed.csv / C_aero.csv / ＊.png … 後処理で生成"],
         Inches(0.4), Inches(4.05), Inches(12.53), Inches(1.15), C_MID_BLUE, C_PROG_BG)

note_bar(slide, "🔁 後処理だけやり直したいとき",
         ["MATLAB で  run_postprocess('C:\\…\\WindyData\\実験名')  の1行。",
          "気温・気圧は実験フォルダの experiment_log.json から自動で読み込まれる。"],
         Inches(0.4), Inches(5.4), Inches(12.53), Inches(1.1), C_GREEN, C_GREEN_LIGHT)


# ============================================================
#  スライド 9: ゼロ揚力角からの原点パルス自動修正
# ============================================================
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "後処理後の対話：ゼロ揚力角の設定修正",
             "実験データから迎角0°の原点パルス(origin_pulse)を求め直し、y/n で更新")
tag_box(slide, "👤 y/n で確認", Inches(11.0), Inches(0.25), C_GREEN, Inches(2.0))

# 仕組みフロー（3ステップ）
flow2 = [
    ("① α₀ を推定", "calc_force.py が線形域から\nゼロ揚力角 α₀ を算出", C_PROG_BG, C_MID_BLUE),
    ("② 推奨値を提示", "推奨 origin_pulse =\n現在値 − round(α₀ × 250)", C_PROG_BG, C_MID_BLUE),
    ("③ y/n で更新", "y なら config.json の\norigin_pulse を自動書換え", C_GREEN_LIGHT, C_GREEN),
]
bw2, bh2, gap2 = Inches(3.7), Inches(1.5), Inches(0.6)
total2 = len(flow2) * bw2 + (len(flow2) - 1) * gap2
sx2 = (SW - total2) / 2
y2 = Inches(1.4)
for i, (ttl, body, bg, fg) in enumerate(flow2):
    x = sx2 + i * (bw2 + gap2)
    add_rect(slide, x, y2, bw2, bh2, fill=bg, border=fg, border_w=Pt(2.2))
    add_text(slide, ttl, x, y2 + Inches(0.1), bw2, Inches(0.4),
             font_size=Pt(15), bold=True, color=fg, align=PP_ALIGN.CENTER)
    add_text(slide, body, x + Inches(0.1), y2 + Inches(0.6), bw2 - Inches(0.2), Inches(0.8),
             font_size=Pt(12), color=C_TEXT, align=PP_ALIGN.CENTER)
    if i < len(flow2) - 1:
        add_text(slide, "▶", x + bw2, y2 + Inches(0.5), gap2, Inches(0.5),
                 font_size=Pt(22), bold=True, color=C_MID_BLUE, align=PP_ALIGN.CENTER)

# 実際の表示例（端末風）
add_rect(slide, Inches(0.9), Inches(3.2), Inches(7.4), Inches(2.4),
         fill=RGBColor(0x1E, 0x1E, 0x1E), border=RGBColor(0x44, 0x44, 0x44))
add_text(slide,
         "==== ゼロ揚力角からの原点パルス修正 ====\n"
         "  推定ゼロ揚力角 α₀ : +0.772°\n"
         "  現在の原点パルス  : 11025 pulse\n"
         "  推奨の原点パルス  : 10832 pulse  (補正 +193 pulse)\n"
         "  ゼロ揚力角の設定(origin_pulse)を\n"
         "  この推奨値に修正しますか？ [y/N]:",
         Inches(1.1), Inches(3.35), Inches(7.0), Inches(2.1),
         font_size=Pt(12.5), color=RGBColor(0x7F, 0xD9, 0x7F))

# 右側の補足
note_bar(slide, "💡 ポイント",
         ["・origin_pulse は config.json の1か所で管理",
          "  （QT_ADL1.m と calc_force.py が共有）。",
          "・計測時の原点は experiment_log に記録され、",
          "  過去実験の再処理でも正しい推奨値になる。",
          "・y を選ぶと次回の実験から新しい原点で計測。",
          "・n なら現状維持。既に一致ならスキップ。",
          "・対称翼で α₀≈0 になるのが理想。"],
         Inches(8.5), Inches(3.2), Inches(4.4), Inches(2.4), C_GREEN, C_GREEN_LIGHT)

note_bar(slide, "⚠ 注意",
         ["α₀（ゼロ揚力角）は翼の取り付け角に対応する。キャンバー翼など α₀≠0 が物理的に正しい",
          "場合もあるため、推奨値の適用は計測者が y/n で判断する（常に従う必要はない）。"],
         Inches(0.4), Inches(5.85), Inches(12.53), Inches(1.05), C_ORANGE, C_ORANGE_LIGHT)


# ============================================================
#  スライド 10: よく使うコマンド早見表
# ============================================================
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "よく使うコマンド早見表", "MATLAB コマンドウィンドウで実行")

cmds = [
    ("run_experiment", C_GREEN, C_GREEN_LIGHT,
     "実験を最初から実行（計測〜後処理〜比較まで）",
     ["画面の指示に従うだけ：",
      "フォルダ名 → 気温・気圧 → 迎角範囲 →",
      "各フェーズでブロワー確認 → 自動計測 → 自動後処理"]),
    ("run_postprocess('…実験フォルダ')", C_MID_BLUE, C_PROG_BG,
     "後処理だけを単体で（再）実行する",
     ["後処理が失敗したときのやり直し、",
      "過去実験の再処理（グラフ範囲変更後など）。",
      "気温・気圧はログから自動取得。"]),
    ("setup_paths", C_PURPLE, RGBColor(0xEE, 0xE7, 0xF6),
     "診断ツール・ヘルパを単体で使う前の準備",
     ["measurement_control / diagnostics を",
      "パスに追加。実行後は QT_ADL1_check_connection、",
      "weight_check, tare_measure 等が使える。"]),
]
y = Inches(1.3)
for cmd, col, bg, sub, lines in cmds:
    h = Inches(1.7)
    add_rect(slide, Inches(0.4), y, Inches(12.53), h, fill=bg, border=col, border_w=Pt(1.8))
    add_rect(slide, Inches(0.4), y, Inches(0.1), h, fill=col)
    # コマンド（端末風）
    add_rect(slide, Inches(0.65), y + Inches(0.2), Inches(5.6), Inches(0.6),
             fill=RGBColor(0x1E, 0x1E, 0x1E))
    add_text(slide, ">> " + cmd, Inches(0.75), y + Inches(0.28), Inches(5.4), Inches(0.45),
             font_size=Pt(14), bold=True, color=RGBColor(0x7F, 0xD9, 0x7F))
    add_text(slide, sub, Inches(0.7), y + Inches(0.95), Inches(5.5), Inches(0.6),
             font_size=Pt(12.5), bold=True, color=col)
    add_text(slide, '\n'.join(lines), Inches(6.5), y + Inches(0.18), Inches(6.3), h - Inches(0.3),
             font_size=Pt(12.5), color=C_TEXT)
    y += h + Inches(0.18)

add_text(slide, "Windy 風洞実験自動計測システム  —  フォルダ構成ガイド",
         Inches(0.4), Inches(7.08), Inches(12.53), Inches(0.32),
         font_size=Pt(10), color=C_GRAYTXT, align=PP_ALIGN.CENTER)


# ============================================================
#  保存
# ============================================================
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Windy_操作マニュアル.pptx")
prs.save(out_path)
print(f"保存完了: {out_path}")
