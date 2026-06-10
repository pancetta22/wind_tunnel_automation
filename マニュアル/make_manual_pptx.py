#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Windy 操作マニュアル PowerPoint 生成スクリプト"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ============================================================
#  カラーパレット
# ============================================================
C_DARK_BLUE   = RGBColor(0x1A, 0x37, 0x6C)   # タイトルバー背景
C_MID_BLUE    = RGBColor(0x2E, 0x6E, 0xB4)   # アクセント
C_LIGHT_BLUE  = RGBColor(0xD6, 0xE8, 0xF7)   # ユーザー操作BOX背景
C_GREEN       = RGBColor(0x21, 0x7A, 0x3C)   # ユーザー操作ラベル
C_GREEN_LIGHT = RGBColor(0xD6, 0xEE, 0xDF)   # ユーザー操作BOX背景
C_ORANGE      = RGBColor(0xD9, 0x6D, 0x00)   # 注意
C_ORANGE_LIGHT= RGBColor(0xFD, 0xF0, 0xDE)   # 注意BOX背景
C_GRAY_BG     = RGBColor(0xF4, 0xF6, 0xF9)   # スライド背景
C_PROG_BG     = RGBColor(0xE8, 0xF0, 0xFB)   # プログラム操作BOX背景
C_PROG_LABEL  = RGBColor(0x2E, 0x6E, 0xB4)   # プログラム操作ラベル
C_TEXT        = RGBColor(0x1A, 0x1A, 0x2E)   # 本文
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_BORDER_G    = RGBColor(0x21, 0x7A, 0x3C)
C_BORDER_B    = RGBColor(0x2E, 0x6E, 0xB4)
C_RED_LIGHT   = RGBColor(0xFD, 0xE8, 0xE8)
C_RED         = RGBColor(0xC0, 0x39, 0x2B)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK_LAYOUT = prs.slide_layouts[6]  # 完全ブランク


# ============================================================
#  ヘルパー関数
# ============================================================

def add_rect(slide, l, t, w, h, fill=None, border=None, border_w=Pt(1)):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = border_w
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             font_size=Pt(14), bold=False, color=C_TEXT,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = font_size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb

def add_label_box(slide, label, label_color, box_color, border_color,
                  content_lines, l, t, w, h,
                  label_size=Pt(11), content_size=Pt(13)):
    """ラベル付きカラーBOX"""
    add_rect(slide, l, t, w, h, fill=box_color, border=border_color, border_w=Pt(1.5))
    # ラベル
    add_text(slide, label, l + Inches(0.1), t + Inches(0.05),
             w - Inches(0.2), Inches(0.28),
             font_size=label_size, bold=True, color=label_color)
    # 本文
    tb = slide.shapes.add_textbox(l + Inches(0.15), t + Inches(0.32),
                                  w - Inches(0.3), h - Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in content_lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = line
        run.font.size  = content_size
        run.font.color.rgb = C_TEXT

def slide_header(slide, title, subtitle=None):
    """共通ヘッダーバー"""
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.85), fill=C_DARK_BLUE)
    add_text(slide, title,
             Inches(0.3), Inches(0.06), Inches(11), Inches(0.6),
             font_size=Pt(24), bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.3), Inches(0.55), Inches(11), Inches(0.35),
                 font_size=Pt(12), color=RGBColor(0xA8, 0xC8, 0xF0))
    # ページ背景
    add_rect(slide, 0, Inches(0.85), prs.slide_width,
             prs.slide_height - Inches(0.85), fill=C_GRAY_BG)

def user_badge(slide, l, t):
    add_rect(slide, l, t, Inches(1.5), Inches(0.32),
             fill=C_GREEN, border=None)
    add_text(slide, "👤 ユーザー操作", l + Inches(0.05), t + Inches(0.02),
             Inches(1.4), Inches(0.28), font_size=Pt(10), bold=True, color=C_WHITE)

def prog_badge(slide, l, t):
    add_rect(slide, l, t, Inches(1.7), Inches(0.32),
             fill=C_MID_BLUE, border=None)
    add_text(slide, "⚙ プログラム自動実行", l + Inches(0.05), t + Inches(0.02),
             Inches(1.6), Inches(0.28), font_size=Pt(10), bold=True, color=C_WHITE)


# ============================================================
#  スライド 1: タイトル
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
# 上半分グラデーション風の背景
add_rect(slide, 0, 0, prs.slide_width, Inches(4.8), fill=C_DARK_BLUE)
add_rect(slide, 0, Inches(4.8), prs.slide_width,
         prs.slide_height - Inches(4.8), fill=C_GRAY_BG)

# タイトル
add_text(slide, "Windy 風洞実験\n自動計測システム",
         Inches(1.0), Inches(1.1), Inches(11.3), Inches(2.5),
         font_size=Pt(44), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

add_text(slide, "操作マニュアル",
         Inches(1.0), Inches(3.2), Inches(11.3), Inches(0.7),
         font_size=Pt(22), color=RGBColor(0xA8, 0xC8, 0xF0), align=PP_ALIGN.CENTER)

# サブ情報
add_rect(slide, Inches(2.5), Inches(4.3), Inches(8.3), Inches(0.5),
         fill=C_MID_BLUE, border=None)
add_text(slide, "run_experiment.m  |  MATLAB R2022a 以降",
         Inches(2.5), Inches(4.3), Inches(8.3), Inches(0.5),
         font_size=Pt(14), color=C_WHITE, align=PP_ALIGN.CENTER)

# キャッチコピー
add_text(slide,
         "迎角ステージ制御・6軸センサ計測・デジボル計測を統合自動化",
         Inches(1.0), Inches(5.2), Inches(11.3), Inches(0.6),
         font_size=Pt(16), color=C_TEXT, align=PP_ALIGN.CENTER, italic=True)


# ============================================================
#  スライド 2: 全体フロー（概要）
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
slide_header(slide, "全体フロー", "実験開始から結果グラフ出力まで")

# フロー図
steps = [
    ("①", "事前準備", "config.json 設定\n機器の接続", C_GREEN_LIGHT, C_GREEN),
    ("②", "スクリプト起動", "run_experiment を\nMATLABで実行", C_GREEN_LIGHT, C_GREEN),
    ("③", "実験情報を入力", "フォルダ名・気温気圧\n最大迎角・開始フェーズ", C_GREEN_LIGHT, C_GREEN),
    ("④", "4フェーズ連続計測", "Pofst→Mofst\n→Pdata→Mdata", C_PROG_BG, C_MID_BLUE),
    ("⑤", "自動後処理", "windspeed.csv生成\n空力係数グラフ出力", C_PROG_BG, C_MID_BLUE),
]

box_w = Inches(2.1)
box_h = Inches(1.6)
gap   = Inches(0.22)
total_w = len(steps) * box_w + (len(steps) - 1) * gap
start_x = (prs.slide_width - total_w) / 2
y = Inches(1.55)

for i, (num, title_s, body, bg, fg) in enumerate(steps):
    x = start_x + i * (box_w + gap)
    add_rect(slide, x, y, box_w, box_h, fill=bg,
             border=fg, border_w=Pt(2.0))
    # 番号円
    add_rect(slide, x + Inches(0.05), y + Inches(0.05),
             Inches(0.38), Inches(0.38), fill=fg, border=None)
    add_text(slide, num, x + Inches(0.05), y + Inches(0.05),
             Inches(0.38), Inches(0.38),
             font_size=Pt(11), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title_s,
             x + Inches(0.1), y + Inches(0.45),
             box_w - Inches(0.2), Inches(0.45),
             font_size=Pt(13), bold=True, color=fg, align=PP_ALIGN.CENTER)
    add_text(slide, body,
             x + Inches(0.1), y + Inches(0.88),
             box_w - Inches(0.2), Inches(0.65),
             font_size=Pt(11), color=C_TEXT, align=PP_ALIGN.CENTER)
    # 矢印
    if i < len(steps) - 1:
        ax = x + box_w + Inches(0.04)
        add_text(slide, "▶", ax, y + Inches(0.6),
                 gap - Inches(0.08), Inches(0.4),
                 font_size=Pt(16), color=C_MID_BLUE, align=PP_ALIGN.CENTER)

# 凡例
ly = Inches(3.55)
add_rect(slide, Inches(0.5), ly, Inches(1.55), Inches(0.32),
         fill=C_GREEN_LIGHT, border=C_GREEN, border_w=Pt(1.5))
add_text(slide, "  👤 ユーザー操作", Inches(0.5), ly + Inches(0.03),
         Inches(1.55), Inches(0.28), font_size=Pt(11), color=C_GREEN, bold=True)
add_rect(slide, Inches(2.3), ly, Inches(1.75), Inches(0.32),
         fill=C_PROG_BG, border=C_MID_BLUE, border_w=Pt(1.5))
add_text(slide, "  ⚙ プログラム自動実行", Inches(2.3), ly + Inches(0.03),
         Inches(1.75), Inches(0.28), font_size=Pt(11), color=C_MID_BLUE, bold=True)

# フェーズ説明
py = Inches(4.05)
add_rect(slide, Inches(0.4), py, prs.slide_width - Inches(0.8), Inches(2.8),
         fill=C_PROG_BG, border=C_MID_BLUE, border_w=Pt(1.5))
add_text(slide, "⚙ 4フェーズとは",
         Inches(0.6), py + Inches(0.1), Inches(4), Inches(0.36),
         font_size=Pt(13), bold=True, color=C_MID_BLUE)

phase_info = [
    ("Pofst", "正迎角・無風",  "0°→+1°→0°→…→+max°→0°", "ゼロ点補正データ（風なし）"),
    ("Mofst", "負迎角・無風",  "0°→-1°→0°→…→-max°→0°", "ゼロ点補正データ（風なし）"),
    ("Pdata", "正迎角・有風",  "0°→+1°→0°→…→+max°→0°", "実験データ（風あり）"),
    ("Mdata", "負迎角・有風",  "0°→-1°→0°→…→-max°→0°", "実験データ（風あり）"),
]
col_w = (prs.slide_width - Inches(0.9)) / 4
for i, (ph, sub, seq, note) in enumerate(phase_info):
    cx = Inches(0.45) + i * col_w
    add_text(slide, ph,
             cx + Inches(0.1), py + Inches(0.5),
             col_w - Inches(0.2), Inches(0.38),
             font_size=Pt(15), bold=True, color=C_DARK_BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, sub,
             cx + Inches(0.1), py + Inches(0.85),
             col_w - Inches(0.2), Inches(0.3),
             font_size=Pt(11), color=C_ORANGE, align=PP_ALIGN.CENTER, bold=True)
    add_text(slide, seq,
             cx + Inches(0.1), py + Inches(1.15),
             col_w - Inches(0.2), Inches(0.38),
             font_size=Pt(10), color=C_TEXT, align=PP_ALIGN.CENTER)
    add_text(slide, note,
             cx + Inches(0.1), py + Inches(1.55),
             col_w - Inches(0.2), Inches(0.38),
             font_size=Pt(10), color=RGBColor(0x55, 0x55, 0x55),
             align=PP_ALIGN.CENTER, italic=True)
    if i < 3:
        add_text(slide, "│", cx + col_w - Inches(0.12),
                 py + Inches(0.9), Inches(0.22), Inches(0.7),
                 font_size=Pt(18), color=RGBColor(0xBB, 0xCC, 0xDD),
                 align=PP_ALIGN.CENTER)


# ============================================================
#  スライド 3: 事前準備（ユーザー）
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
slide_header(slide, "① 事前準備", "実験前に1回だけ確認する設定")
user_badge(slide, Inches(10.8), Inches(0.25))

y0 = Inches(1.05)

# LEFT: config.json
add_label_box(slide, "📄 config.json の設定",
              C_GREEN, C_GREEN_LIGHT, C_GREEN,
              [
                  'output_dir    : データ保存先のルートフォルダ',
                  '                例) "C:/Users/xxx/WindyData"',
                  'leptrino_port : Leptrino の COM ポート番号',
                  'qt_adl1_port  : 迎角ステージの COM ポート',
                  'r6441b_port   : デジボルの COM ポート',
                  'python_exe    : 32bit Python（計測用）のパス',
                  'python_exe_64 : 64bit Python（後処理用）のパス',
                  '',
                  'calib_a / calib_b : 較正定数（通常変更不要）',
              ],
              Inches(0.4), y0, Inches(5.7), Inches(4.0),
              content_size=Pt(12))

# RIGHT: 機器接続
add_label_box(slide, "🔌 機器の接続確認",
              C_GREEN, C_GREEN_LIGHT, C_GREEN,
              [
                  '① Leptrino 6軸センサ',
                  '    USB-RS232C 変換経由で PC に接続',
                  '',
                  '② 迎角ステージ（QT-ADL1）',
                  '    USB-RS232C 変換経由で PC に接続',
                  '',
                  '③ R6441B デジボル（差圧センサ）',
                  '    RS-232C で PC に接続',
              ],
              Inches(6.4), y0, Inches(6.5), Inches(4.0),
              content_size=Pt(12))

# 注意BOX
add_label_box(slide, "⚠  注意",
              C_ORANGE, C_ORANGE_LIGHT, C_ORANGE,
              [
                  'config.json の output_dir は "WindyData" のルートフォルダを指定（日付サブフォルダは含めない）',
                  '実験ごとに自動でサブフォルダが作成されます',
              ],
              Inches(0.4), Inches(5.2), Inches(12.5), Inches(0.9),
              content_size=Pt(12))


# ============================================================
#  スライド 4: スクリプト起動 + 入力ステップ
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
slide_header(slide, "② スクリプト起動 ＆ 実験情報の入力", "MATLAB コマンドウィンドウで操作")
user_badge(slide, Inches(10.8), Inches(0.25))

# 起動コマンド
add_rect(slide, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.8),
         fill=C_DARK_BLUE, border=None)
add_text(slide, ">> run_experiment",
         Inches(0.7), Inches(1.15), Inches(10), Inches(0.55),
         font_size=Pt(20), bold=True, color=RGBColor(0x7F, 0xD9, 0x7F))

inputs = [
    ("実験フォルダ名",
     "実験フォルダ名: ",
     "260605_rigid",
     [
         "WindyData の直下に作成されるフォルダ名",
         "形式の例: 260605_rigid  /  260605_flex",
         "使用不可文字:  / \\ : * ? \" < > |",
     ]),
    ("気温・気圧",
     "気温 [℃]: \n気圧 [mmHg]: ",
     "25.3\n758.0",
     [
         "空気密度 ρ が自動計算される",
         "水密度もKellの式で自動計算",
         "全フェーズで共通（1回のみ入力）",
     ]),
    ("最大迎角",
     "最大迎角 [度, 1-30]: ",
     "30 ← Enter のみで 30°",
     [
         "1 〜 30 の整数で入力",
         "入力した角度まで 1°刻みで計測",
         "1フェーズあたり 2×max+1 点",
     ]),
    ("開始フェーズ",
     "開始フェーズ [1-4]: ",
     "1 ← Enter のみで Pofst から",
     [
         "通常は Enter（=1: Pofst から）",
         "途中再開時は 2〜4 を選択",
         "1:Pofst  2:Mofst  3:Pdata  4:Mdata",
     ]),
]

col_w = Inches(3.05)
gap_x = Inches(0.08)
start_x = Inches(0.4)
y1 = Inches(2.0)
bh = Inches(3.9)

for i, (label, prompt, response, notes) in enumerate(inputs):
    x = start_x + i * (col_w + gap_x)

    # 番号
    add_rect(slide, x, y1, Inches(0.32), Inches(0.32),
             fill=C_MID_BLUE, border=None)
    add_text(slide, str(i+1), x, y1, Inches(0.32), Inches(0.32),
             font_size=Pt(12), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # ラベル
    add_text(slide, label, x + Inches(0.36), y1 + Inches(0.04),
             col_w - Inches(0.36), Inches(0.3),
             font_size=Pt(13), bold=True, color=C_DARK_BLUE)

    # 端末っぽいBOX
    add_rect(slide, x, y1 + Inches(0.38), col_w, Inches(1.05),
             fill=RGBColor(0x1E, 0x1E, 0x1E), border=RGBColor(0x44, 0x44, 0x44))
    add_text(slide, prompt, x + Inches(0.1), y1 + Inches(0.42),
             col_w - Inches(0.2), Inches(0.42),
             font_size=Pt(10), color=RGBColor(0xCC, 0xCC, 0xCC), italic=True)
    add_text(slide, response, x + Inches(0.1), y1 + Inches(0.82),
             col_w - Inches(0.2), Inches(0.52),
             font_size=Pt(10), bold=True, color=RGBColor(0x7F, 0xD9, 0x7F))

    # 説明
    add_rect(slide, x, y1 + Inches(1.48), col_w, bh - Inches(1.48),
             fill=C_GREEN_LIGHT, border=C_GREEN, border_w=Pt(1.2))
    for j, note in enumerate(notes):
        add_text(slide, "• " + note,
                 x + Inches(0.1), y1 + Inches(1.55) + j * Inches(0.7),
                 col_w - Inches(0.2), Inches(0.65),
                 font_size=Pt(11), color=C_TEXT)


# ============================================================
#  スライド 5: プログラムが自動で行うこと（機器接続・フォルダ作成）
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
slide_header(slide, "入力後にプログラムが自動で行うこと", "機器接続・フォルダ作成")
prog_badge(slide, Inches(10.6), Inches(0.25))

y0 = Inches(1.05)

add_label_box(slide, "⚙ 自動処理の流れ",
              C_MID_BLUE, C_PROG_BG, C_MID_BLUE,
              [
                  "1.  迎角ステージ（QT-ADL1）に接続  →  ホームポジションへ自動復帰",
                  "2.  Leptrino 6軸センサの接続を確認",
                  "3.  R6441B デジボルに接続",
                  "4.  WindyMonitor（モニター画面）を起動",
                  "5.  experiment_log.json を保存（気温・気圧・空気密度など）",
              ],
              Inches(0.4), y0, Inches(12.5), Inches(2.5),
              content_size=Pt(13))

# フォルダ構造
add_label_box(slide, "📁 自動作成されるフォルダ構造",
              C_MID_BLUE, C_PROG_BG, C_MID_BLUE,
              [], Inches(0.4), Inches(3.7), Inches(12.5), Inches(3.55))

folder_lines = [
    ("WindyData/",                                 0, C_DARK_BLUE,  True,  Pt(14)),
    ("└── 260605_rigid/",                          1, C_MID_BLUE,   True,  Pt(13)),
    ("    ├── data/",                              2, C_MID_BLUE,   True,  Pt(13)),
    ("    │   ├── 20260605_143022_260605_Pofst_01.00.csv   ← 6軸センサCSV",  3, C_TEXT, False, Pt(11)),
    ("    │   ├── 20260605_143022_260605_Pofst_01.00_volt_raw.csv",           3, C_TEXT, False, Pt(11)),
    ("    │   └── ...",                            3, C_TEXT,       False, Pt(11)),
    ("    ├── 20260605_Pofst_volt_summary.csv      ← 差圧電圧サマリー（フェーズ毎）",  2, C_TEXT, False, Pt(11)),
    ("    ├── 20260605_experiment_log.json          ← 気温・気圧・校正定数",           2, C_TEXT, False, Pt(11)),
    ("    ├── windspeed.csv                         ← 後処理で自動生成",              2, C_TEXT, False, Pt(11)),
    ("    └── Cl.png / Cd.png / ...                ← 後処理で自動生成",              2, C_TEXT, False, Pt(11)),
]

for j, (txt, indent, col, bd, sz) in enumerate(folder_lines):
    add_text(slide, txt,
             Inches(0.7), Inches(3.85) + j * Inches(0.3),
             Inches(12.1), Inches(0.3),
             font_size=sz, bold=bd, color=col)


# ============================================================
#  スライド 6: 各フェーズ開始時のユーザー操作
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
slide_header(slide, "各フェーズ開始時のユーザー操作", "フェーズは Pofst→Mofst→Pdata→Mdata の順に自動進行")
user_badge(slide, Inches(10.8), Inches(0.25))

y0 = Inches(1.05)

# 無風フェーズ
add_label_box(slide, "👤 無風フェーズ（Pofst / Mofst）の確認",
              C_GREEN, C_GREEN_LIGHT, C_GREEN,
              [], Inches(0.4), y0, Inches(6.0), Inches(5.65))

ofst_steps = [
    ("Step 1", "ブロワー停止の確認",
     "「ブロワーが停止していることを\n確認してください。\nEnter を押してください: 」\n\n→ ブロワー停止を確認して Enter\n\n※ Pofst の直前にのみ、差圧センサの\n   電圧オフセットが自動計測されます\n   （ユーザー操作不要）"),
]
for i, (step, ttl, body) in enumerate(ofst_steps):
    sy = y0 + Inches(0.42) + i * Inches(2.45)
    add_rect(slide, Inches(0.55), sy, Inches(0.7), Inches(0.28),
             fill=C_GREEN, border=None)
    add_text(slide, step, Inches(0.55), sy, Inches(0.7), Inches(0.28),
             font_size=Pt(10), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, ttl, Inches(1.3), sy, Inches(4.8), Inches(0.28),
             font_size=Pt(12), bold=True, color=C_DARK_BLUE)
    add_text(slide, body, Inches(0.6), sy + Inches(0.32),
             Inches(5.5), Inches(2.0),
             font_size=Pt(11), color=C_TEXT)

# 有風フェーズ
add_label_box(slide, "👤 有風フェーズ（Pdata / Mdata）の確認",
              C_GREEN, C_GREEN_LIGHT, C_GREEN,
              [], Inches(6.7), y0, Inches(6.2), Inches(5.65))

data_steps = [
    ("Step 1", "ブロワー起動・風速安定待ち",
     "「ブロワーを起動し、風速が安定した\nことを確認してください。\nEnter を押してください: 」\n\n→ ブロワーを起動\n→ 風速が安定するのを確認\n→ Enter を押す"),
]
for i, (step, ttl, body) in enumerate(data_steps):
    sy = y0 + Inches(0.42) + i * Inches(2.45)
    add_rect(slide, Inches(6.85), sy, Inches(0.7), Inches(0.28),
             fill=C_GREEN, border=None)
    add_text(slide, step, Inches(6.85), sy, Inches(0.7), Inches(0.28),
             font_size=Pt(10), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, ttl, Inches(7.6), sy, Inches(5.0), Inches(0.28),
             font_size=Pt(12), bold=True, color=C_DARK_BLUE)
    add_text(slide, body, Inches(6.9), sy + Inches(0.32),
             Inches(5.7), Inches(2.0),
             font_size=Pt(11), color=C_TEXT)

# 注意
add_label_box(slide, "⚠  有風フェーズのブロワー操作は 1回だけ（Pdata開始前のみ）",
              C_ORANGE, C_ORANGE_LIGHT, C_ORANGE,
              ["Mdata はそのまま風を出し続けた状態で自動継続されます。",
               "Mdata 終了後のブロワー停止はユーザーが手動で行ってください。"],
              Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.78),
              content_size=Pt(11))


# ============================================================
#  スライド 7: 1計測点の自動処理フロー
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
slide_header(slide, "1計測点ごとにプログラムが自動で行うこと",
             "各迎角で繰り返される処理（例: 0°→+1°→0°→+2°→0°…）")
prog_badge(slide, Inches(10.6), Inches(0.25))

steps7 = [
    ("a", "迎角ステージ移動", "指定角度へステージを自動移動\n位置が確定するまで待機"),
    ("b", "振動収束待ち",      f"ステージ停止後 {'{angle_settle_sec}'} 秒待機\n（config.json の angle_settle_sec で設定）"),
    ("c", "6軸センサ計測開始", "leptrino_server.py を\nバックグラウンドで起動\nCSVへのストリーミング記録を開始"),
    ("d", "デジボル計測ループ", "6軸センサが規定量に達するまで\nR6441Bに繰り返しポーリング\n差圧電圧を収集（~10サンプル/秒）"),
    ("e", "データ保存",         "6軸センサCSV を確認\n差圧電圧を volt_raw.csv に保存\nvolt_summary.csv に1行追記"),
]

box_w2 = Inches(2.35)
box_h2 = Inches(3.5)
start_x2 = Inches(0.4)
gap2 = Inches(0.18)
y2 = Inches(1.1)

for i, (alpha, title_s, body) in enumerate(steps7):
    x = start_x2 + i * (box_w2 + gap2)
    add_rect(slide, x, y2, box_w2, box_h2,
             fill=C_PROG_BG, border=C_MID_BLUE, border_w=Pt(2.0))
    # アルファベット円
    add_rect(slide, x + Inches(0.05), y2 + Inches(0.06),
             Inches(0.4), Inches(0.4), fill=C_MID_BLUE, border=None)
    add_text(slide, alpha, x + Inches(0.05), y2 + Inches(0.06),
             Inches(0.4), Inches(0.4),
             font_size=Pt(14), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title_s,
             x + Inches(0.1), y2 + Inches(0.52),
             box_w2 - Inches(0.2), Inches(0.5),
             font_size=Pt(13), bold=True, color=C_DARK_BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, body,
             x + Inches(0.12), y2 + Inches(1.05),
             box_w2 - Inches(0.24), Inches(2.3),
             font_size=Pt(11), color=C_TEXT, align=PP_ALIGN.CENTER)
    if i < len(steps7) - 1:
        ax = x + box_w2 + Inches(0.03)
        add_text(slide, "▶", ax, y2 + Inches(1.4),
                 gap2 - Inches(0.03), Inches(0.5),
                 font_size=Pt(18), color=C_MID_BLUE, align=PP_ALIGN.CENTER)

# 補足
add_label_box(slide, "⚙ モニター画面（WindyMonitor）の表示内容",
              C_MID_BLUE, C_PROG_BG, C_MID_BLUE,
              [
                  "・現在の迎角 / 進捗バー（何点目/何点中）",
                  "・6軸センサの振動波形（直近0.5秒をオシロスコープ的に表示）",
                  "・差圧電圧のリアルタイム値",
                  "・ファイルサイズ進捗バー（6軸センサ計測の完了度を視覚化）",
                  "・「一時停止」ボタン → 押すと「再開」「停止」の2ボタンが現れる",
              ],
              Inches(0.4), Inches(4.75), Inches(12.5), Inches(2.1),
              content_size=Pt(12))


# ============================================================
#  スライド 8: エラー発生時の対処
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
slide_header(slide, "エラー発生時の対処", "計測中にエラーが起きても対話的に回復できる")
user_badge(slide, Inches(10.8), Inches(0.25))

# エラーBOX
add_rect(slide, Inches(0.4), Inches(1.05), Inches(12.5), Inches(1.1),
         fill=C_RED_LIGHT, border=C_RED, border_w=Pt(2.0))
add_text(slide,
         "⚠  エラー発生時 または「停止」ボタン押下時にプログラムが表示するメッセージ（例）",
         Inches(0.6), Inches(1.08), Inches(12.0), Inches(0.35),
         font_size=Pt(12), bold=True, color=C_RED)
add_text(slide,
         "════════════════════════════════════════\n"
         "  [エラー] 計測点 5/61 でエラーが発生しました\n"
         "  フェーズ: Pdata  内容: Timeout reading from serial port\n"
         "════════════════════════════════════════\n"
         "どうしますか？  R: 再試行  S: スキップ  P: フェーズ再開  C: フェーズ選択  Q: 終了",
         Inches(0.6), Inches(1.42), Inches(12.0), Inches(0.65),
         font_size=Pt(10.5), color=C_TEXT)

# 選択肢一覧
choices = [
    ("R", "再試行", C_MID_BLUE,
     "この計測点をもう一度実行する\n（機器の瞬間的なエラー時に使用）"),
    ("S", "スキップ", RGBColor(0x6A, 0x0A, 0x8A),
     "この点を飛ばして次の計測点へ進む\n（1点欠損でよい場合）"),
    ("P", "フェーズ\n再開", C_ORANGE,
     "このフェーズを最初からやり直す\n部分データは自動削除される"),
    ("C", "フェーズ\n選択", RGBColor(0x80, 0x60, 0x00),
     "やり直すフェーズ番号を選ぶ\n（前のフェーズからやり直す場合）"),
    ("Q", "終了", C_RED,
     "実験を中止してプログラムを終了\n機器の接続は自動でクリーンアップ"),
]

box_w3 = Inches(2.35)
y3 = Inches(2.3)
bh3 = Inches(3.6)
for i, (key, title_s, col, body) in enumerate(choices):
    x = Inches(0.4) + i * (box_w3 + Inches(0.1))
    add_rect(slide, x, y3, box_w3, bh3, fill=C_GRAY_BG, border=col, border_w=Pt(2.5))
    add_rect(slide, x, y3, box_w3, Inches(0.7), fill=col, border=None)
    add_text(slide, key, x, y3, Inches(0.7), Inches(0.7),
             font_size=Pt(28), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title_s, x + Inches(0.72), y3 + Inches(0.12),
             box_w3 - Inches(0.78), Inches(0.5),
             font_size=Pt(14), bold=True, color=C_WHITE)
    add_text(slide, body, x + Inches(0.12), y3 + Inches(0.82),
             box_w3 - Inches(0.24), bh3 - Inches(0.95),
             font_size=Pt(12), color=C_TEXT)

# 注意
add_label_box(slide, "💡 フェーズ再開（P/C）時の部分データ自動削除について",
              C_ORANGE, C_ORANGE_LIGHT, C_ORANGE,
              ['P または C でフェーズを最初からやり直す際、'
               'data/ フォルダ内の当該フェーズのファイルを自動削除してから再計測します。',
               '後処理（calc_force.py）に余分なデータが混入することを防ぐためです。'],
              Inches(0.4), Inches(6.0), Inches(12.5), Inches(1.0),
              content_size=Pt(11))


# ============================================================
#  スライド 9: 全フェーズ完了後の後処理
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
slide_header(slide, "全フェーズ完了後：自動後処理", "make_windspeed.py → calc_force.py を順次実行")
prog_badge(slide, Inches(10.6), Inches(0.25))

y0 = Inches(1.05)

# Step 1
add_label_box(slide, "⚙ Step 1：windspeed.csv の生成（make_windspeed.py）",
              C_MID_BLUE, C_PROG_BG, C_MID_BLUE,
              [
                  "【初回のみ】post_process/venv/ に 64bit Python 仮想環境を自動作成",
                  "          requirements.txt から pandas / numpy / scipy 等を pip install",
                  "差圧電圧サマリー（Pdata / Mdata の volt_summary.csv）を読み込む",
                  "experiment_log.json から気温・気圧・校正定数を取得",
                  "各計測点の差圧電圧 → 動圧水柱高さ → 風速 U [m/s] を計算",
                  "windspeed.csv として実験フォルダに保存",
              ],
              Inches(0.4), y0, Inches(12.5), Inches(2.0),
              content_size=Pt(13))

# Step 2
add_label_box(slide, "⚙ Step 2：空力係数グラフの生成（calc_force.py）",
              C_MID_BLUE, C_PROG_BG, C_MID_BLUE,
              [
                  "windspeed.csv と data/ 内の 6軸センサ CSV を読み込む",
                  "オフセット補正（Pofst / Mofst データを差し引き）を実施",
                  "Cl, Cd, Cm などの空力係数を迎角ごとに計算",
                  "グラフ（PNG）を実験フォルダに自動保存",
              ],
              Inches(0.4), Inches(3.2), Inches(12.5), Inches(2.0),
              content_size=Pt(13))

# 完了後のユーザー確認
add_label_box(slide, "👤 後処理完了後にユーザーがすること",
              C_GREEN, C_GREEN_LIGHT, C_GREEN,
              [
                  "① WindyData/<実験フォルダ>/ を開いて Cl.png / Cd.png などを確認",
                  "② ブロワーを停止する（有風実験の場合）",
                  "③ 迎角ステージは自動でホームポジション（0°）に戻っています",
              ],
              Inches(0.4), Inches(5.35), Inches(12.5), Inches(1.65),
              content_size=Pt(13))

# 後処理がスキップされる条件
add_label_box(slide, "⚠  後処理がスキップされるケース",
              C_ORANGE, C_ORANGE_LIGHT, C_ORANGE,
              [
                  "4フェーズのいずれか volt_summary.csv が欠けている場合は後処理を自動スキップ",
                  "→ 欠けているフェーズを再実行するか、手動で後処理スクリプトを実行してください",
              ],
              Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.85),
              content_size=Pt(11))


# ============================================================
#  スライド 10: まとめ（チートシート）
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
slide_header(slide, "まとめ：操作チートシート", "実験当日に確認するポイント")

y0 = Inches(1.05)

add_label_box(slide, "👤 ユーザーが入力・確認すること（全部で 7ステップ）",
              C_GREEN, C_GREEN_LIGHT, C_GREEN,
              [
                  "1.  run_experiment を実行",
                  "2.  実験フォルダ名を入力   例: 260605_rigid",
                  "3.  気温 [℃] を入力",
                  "4.  気圧 [mmHg] を入力",
                  "5.  最大迎角を入力   （Enter のみ → 30°）",
                  "6.  開始フェーズを選択   （Enter のみ → 1: Pofst から）",
                  "--- フェーズごとに ---",
                  "7a. 無風フェーズ → ブロワー停止確認 → Enter",
                  "    （Pofst 直前のみ電圧オフセットが自動計測される、操作不要）",
                  "7b. 有風フェーズ → ブロワー起動・風速安定後 → Enter",
              ],
              Inches(0.4), y0, Inches(6.2), Inches(5.75),
              content_size=Pt(12))

add_label_box(slide, "⚙ プログラムが自動でやること",
              C_MID_BLUE, C_PROG_BG, C_MID_BLUE,
              [
                  "・機器（ステージ/センサ/デジボル）の接続",
                  "・フォルダ自動作成（WindyData/<実験名>/data/）",
                  "・各計測点：ステージ移動→センサ計測→デジボル計測→保存",
                  "・4フェーズを Pofst→Mofst→Pdata→Mdata の順に自動進行",
                  "・エラー時：R（再試行）/ S（スキップ）/ P（フェーズ再開）…を確認",
                  "・全完了後：windspeed.csv → 空力係数グラフを自動生成",
              ],
              Inches(6.7), y0, Inches(6.2), Inches(5.75),
              content_size=Pt(12))

# フッターライン
add_rect(slide, Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.03),
         fill=C_MID_BLUE, border=None)
add_text(slide, "Windy 風洞実験自動計測システム  —  操作マニュアル",
         Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.35),
         font_size=Pt(10), color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)


# ============================================================
#  保存
# ============================================================
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Windy_操作マニュアル.pptx")
prs.save(out_path)
print(f"保存完了: {out_path}")
