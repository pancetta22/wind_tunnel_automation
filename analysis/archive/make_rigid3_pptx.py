#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""rigid3（デジボル接触不良 修正後）結果まとめ — 揚力傾斜の謎の解決確認。±5°。"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

C_DARK=RGBColor(0x1A,0x37,0x6C); C_MID=RGBColor(0x2E,0x6E,0xB4)
C_GREEN=RGBColor(0x21,0x7A,0x3C); C_GREEN_L=RGBColor(0xD6,0xEE,0xDF)
C_ORANGE=RGBColor(0xD9,0x6D,0x00); C_ORANGE_L=RGBColor(0xFD,0xF0,0xDE)
C_GRAY=RGBColor(0xF4,0xF6,0xF9); C_TEXT=RGBColor(0x1A,0x1A,0x2E)
C_WHITE=RGBColor(0xFF,0xFF,0xFF); C_RED=RGBColor(0xC0,0x39,0x2B); C_SUB=RGBColor(0x5A,0x6A,0x85)

prs=Presentation(); prs.slide_width=Inches(13.33); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]

def rect(s,l,t,w,h,fill=None,border=None,bw=Pt(1)):
    sh=s.shapes.add_shape(1,l,t,w,h); sh.shadow.inherit=False
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    else: sh.fill.background()
    if border: sh.line.color.rgb=border; sh.line.width=bw
    else: sh.line.fill.background()
    return sh

def text(s,t,l,tp,w,h,size=Pt(14),bold=False,color=C_TEXT,align=PP_ALIGN.LEFT,it=False):
    tb=s.shapes.add_textbox(l,tp,w,h); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=t
    r.font.size=size; r.font.bold=bold; r.font.color.rgb=color; r.font.italic=it
    return tb

def header(s,title,sub=None):
    rect(s,0,Inches(0.85),prs.slide_width,prs.slide_height-Inches(0.85),fill=C_GRAY)
    rect(s,0,0,prs.slide_width,Inches(0.85),fill=C_DARK)
    text(s,title,Inches(0.35),Inches(0.08),Inches(12.5),Inches(0.55),size=Pt(23),bold=True,color=C_WHITE)
    if sub: text(s,sub,Inches(0.37),Inches(0.56),Inches(12),Inches(0.3),size=Pt(11.5),color=RGBColor(0xA8,0xC8,0xF0))

def bullets(s,lines,l,t,w,h,size=Pt(13),gap=Pt(6)):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True; first=True
    for it in lines:
        txt=it[0]; lv=it[1] if len(it)>1 else 0; col=it[2] if len(it)>2 else C_TEXT; bd=it[3] if len(it)>3 else False
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.space_before=gap; p.level=lv
        mark="■ " if lv==0 else "・ "
        r=p.add_run(); r.text=mark+txt; r.font.size=size if lv==0 else Pt(size.pt-1.5)
        r.font.color.rgb=col; r.font.bold=bd

def panel(s,l,t,w,h,fill,border): rect(s,l,t,w,h,fill=fill,border=border,bw=Pt(1.5))

# ===== S1 サマリ（解決確認）=====
s=prs.slides.add_slide(BLANK)
header(s,"【解決】揚力傾斜の系統差 = デジボル接触不良が真因","rigid3（接触修正後）で旧データと一致。揚力傾斜は ±5° で算出")

for (l,title,val,unit,sub) in [
 (0.4,"揚力傾斜 dCl/dα","3.87","/rad","旧3.94 と一致（0.98×）"),
 (4.6,"ゼロ揚力角 α₀","0.77","°","旧0.5〜0.7 と一致"),
 (8.8,"風速 U","13.11","m/s","旧12.2より高風速でも傾斜同じ")]:
    panel(s,Inches(l),Inches(1.1),Inches(4.0 if l<8 else 4.1),Inches(2.5),C_GREEN_L,C_GREEN)
    text(s,title,Inches(l+0.15),Inches(1.2),Inches(3.7),Inches(0.4),size=Pt(13),bold=True,color=C_GREEN)
    text(s,val,Inches(l+0.15),Inches(1.65),Inches(3.7),Inches(0.9),size=Pt(40),bold=True,color=C_DARK)
    text(s,unit,Inches(l+0.15),Inches(2.62),Inches(3.7),Inches(0.4),size=Pt(13),color=C_SUB)
    text(s,sub,Inches(l+0.15),Inches(2.98),Inches(3.8),Inches(0.5),size=Pt(11),color=C_SUB)

panel(s,Inches(0.4),Inches(3.8),Inches(12.5),Inches(3.05),C_GREEN_L,C_GREEN)
text(s,"結論",Inches(0.55),Inches(3.9),Inches(12),Inches(0.4),size=Pt(15),bold=True,color=C_GREEN)
bullets(s,[
 ("接触不良を直したら、揚力傾斜 3.87/rad ＝ 旧3.94/rad とほぼ完全一致。1.14〜1.29×の差は消滅",0,C_GREEN,True),
 ("rigid3 は U=13.11 m/s（旧より高風速）でも傾斜は同じ → 同じ翼なら傾斜は風速によらない（無次元化が正しく機能）",0),
 ("真犯人：差圧電圧（デジボル）の接触不良 → mV過小 → 動圧q過小 → Cl=L/q が過大 → 傾斜が過大に見えていた",0,C_RED,True),
 ("これまで疑った マウント/センサFy軸/定格/処理/空力弾性 は すべて無罪",0,C_SUB),
],Inches(0.55),Inches(4.35),Inches(12.1),Inches(2.4),size=Pt(13.5),gap=Pt(9))

# ===== S2 推移表 =====
s=prs.slides.add_slide(BLANK)
header(s,"推移：605 → 608 → 610 → rigid3（剛体 NACA0012, ±5°）","記録mVが低い＝接触不良が酷い ほど 傾斜が過大に出ていた")
data=[
 ("旧 平均","1165","3.94","—","基準（正常）",False),
 ("新0605","1054","4.92","→3.90","接触最悪→mV最低→傾斜最大",False),
 ("新0608","1169","4.44","→3.92","",False),
 ("新0610","1157","4.49","→3.92","",False),
 ("新 rigid3 ★","1332","3.87","実測","接触修正後→補正不要で旧に一致",True),
]
heads=["実験","記録mV","記録 dCl/dα","真mV1320へ補正","備考"]
xs=[0.5,2.6,4.0,6.0,8.3]; ws=[2.1,1.4,2.0,2.3,4.7]
y=Inches(1.25)
rect(s,Inches(0.4),y,Inches(12.5),Inches(0.42),fill=C_DARK)
for hx,hw,ht in zip(xs,ws,heads):
    text(s,ht,Inches(hx),y+Inches(0.04),Inches(hw),Inches(0.34),size=Pt(12.5),bold=True,color=C_WHITE)
y=Inches(1.71)
for i,(lab,mv,rec,corr,note,hl) in enumerate(data):
    bg=C_GREEN_L if hl else (C_WHITE if i%2==0 else RGBColor(0xEC,0xF1,0xF8))
    rect(s,Inches(0.4),y,Inches(12.5),Inches(0.62),fill=bg,border=(C_GREEN if hl else RGBColor(0xD0,0xDA,0xE8)),bw=Pt(2 if hl else 0.75))
    col=C_GREEN if hl else C_TEXT
    for hx,hw,v in zip(xs,ws,[lab,mv,rec,corr,note]):
        text(s,v,Inches(hx),y+Inches(0.13),Inches(hw),Inches(0.4),size=Pt(12.5),bold=hl,color=col)
    y+=Inches(0.62)
panel(s,Inches(0.4),Inches(5.4),Inches(12.5),Inches(1.45),C_ORANGE_L,C_ORANGE)
text(s,"読み方",Inches(0.55),Inches(5.5),Inches(12),Inches(0.4),size=Pt(13),bold=True,color=C_ORANGE)
bullets(s,[
 ("0605は記録mVが1054と最低（接触が最悪）→ 傾斜4.92と最大。接触の程度がそのまま傾斜の過大さに対応",0),
 ("真mV=1320へ補正すると 0605/0608/0610 すべて 約3.9/rad に収束。rigid3は実測で既に3.87（補正不要）",0,C_GREEN,True),
],Inches(0.55),Inches(5.9),Inches(12.1),Inches(0.9),size=Pt(12.5),gap=Pt(5))

# ===== S3 真因と二つの別問題 =====
s=prs.slides.add_slide(BLANK)
header(s,"真因の整理：『傾斜』と『α₀』は別原因だった")
panel(s,Inches(0.4),Inches(1.1),Inches(6.15),Inches(3.6),C_RED_LIGHT if False else C_WHITE,C_RED)
text(s,"① 揚力傾斜 1.14〜1.29× → デジボル接触不良",Inches(0.55),Inches(1.2),Inches(5.85),Inches(0.5),size=Pt(14),bold=True,color=C_RED)
bullets(s,[
 ("差圧電圧の接触不良で mV が過小に読めていた",0,C_RED,True),
 ("mV低 → 風速U低 → 動圧q低",0),
 ("Cl = L / q なので q過小 → Cl・傾斜が過大",0),
 ("接触修正(rigid3, mV=1332)で 3.87/rad＝旧に一致",0,C_GREEN,True),
 ("※傾斜は風速によらない量。同じ翼なら一定のはずが、",1,C_SUB),
 ("  qの誤りで膨らんで見えていた",1,C_SUB),
],Inches(0.55),Inches(1.7),Inches(5.85),Inches(2.9),size=Pt(12.5),gap=Pt(6))

panel(s,Inches(6.75),Inches(1.1),Inches(6.15),Inches(3.6),C_WHITE,C_MID)
text(s,"② ゼロ揚力角 α₀ のズレ → 取付角（マウント）",Inches(6.9),Inches(1.2),Inches(5.85),Inches(0.5),size=Pt(14),bold=True,color=C_MID)
bullets(s,[
 ("605: α₀=1.44° → 608:0.93 → 610:0.75 → rigid3:0.77",0),
 ("再マウントで取付角が改善し旧（0.5〜0.7）へ収束",0,C_GREEN,True),
 ("α₀は q補正の影響を受けない（Cl=0の位置は不変）",0),
 ("→ デジボルとは独立した別問題だった",0,C_RED,True),
],Inches(6.9),Inches(1.7),Inches(5.85),Inches(2.9),size=Pt(12.5),gap=Pt(7))

panel(s,Inches(0.4),Inches(4.9),Inches(12.5),Inches(1.9),C_GREEN_L,C_GREEN)
text(s,"潰してきた容疑者（すべて無罪と確認済み）",Inches(0.55),Inches(5.0),Inches(12),Inches(0.4),size=Pt(13),bold=True,color=C_GREEN)
bullets(s,[
 ("ソフト/ロガー（瞬時値一致）・センサFz較正(おもり0.25%)・定格Limit・後処理(新旧同一)・整定時間・三角関数",0),
 ("マウントねじれ柔性（空力弾性）も不要に → 真因はデジボル接触不良で確定",0,C_RED,True),
],Inches(0.55),Inches(5.4),Inches(12.1),Inches(1.3),size=Pt(12.5),gap=Pt(6))

# ===== S4 まとめ =====
s=prs.slides.add_slide(BLANK)
rect(s,0,0,prs.slide_width,prs.slide_height,fill=C_DARK)
rect(s,Inches(0.9),Inches(0.9),Inches(0.16),Inches(0.7),fill=C_GREEN)
text(s,"まとめ：謎の解決",Inches(1.2),Inches(0.85),Inches(11),Inches(0.8),size=Pt(30),bold=True,color=C_WHITE)

panel(s,Inches(0.9),Inches(2.0),Inches(11.5),Inches(1.5),RGBColor(0x14,0x2A,0x55),C_GREEN)
bullets(s,[
 ("揚力傾斜の系統差（最大1.29×）の真因 = 差圧デジボルの接触不良（mV過小→q過小→Cl過大）",0,C_WHITE,True),
 ("接触修正後の rigid3 で 3.87/rad ＝ 旧3.94/rad と一致 → 解決",0,RGBColor(0xC8,0xF0,0xD8)),
],Inches(1.15),Inches(2.15),Inches(11.0),Inches(1.2),size=Pt(13.5),gap=Pt(8))

panel(s,Inches(0.9),Inches(3.7),Inches(11.5),Inches(1.5),RGBColor(0x14,0x2A,0x55),C_MID)
bullets(s,[
 ("α₀のズレは別問題（取付角）で、再マウントにより既に解消",0,C_WHITE),
 ("マウント/センサ/定格/処理/空力弾性 はすべて無罪と確認済み",0,RGBColor(0xC8,0xD8,0xF0)),
],Inches(1.15),Inches(3.85),Inches(11.0),Inches(1.2),size=Pt(13.5),gap=Pt(8))

panel(s,Inches(0.9),Inches(5.4),Inches(11.5),Inches(1.3),RGBColor(0x10,0x2A,0x18),C_GREEN)
text(s,"今後",Inches(1.15),Inches(5.5),Inches(10),Inches(0.4),size=Pt(15),bold=True,color=C_GREEN)
bullets(s,[
 ("デジボルの接触を点検項目に追加（計測前にmV妥当性チェック）。過去の新データは真mVが分かれば補正可能",0,C_WHITE,True),
],Inches(1.15),Inches(5.9),Inches(11.0),Inches(0.7),size=Pt(13),gap=Pt(6))

OUT=os.path.join(os.path.dirname(os.path.abspath(__file__)),"Windy_rigid3_解決確認.pptx")
prs.save(OUT)
print("生成完了:",OUT,"(%dスライド)"%len(prs.slides._sldIdLst))
