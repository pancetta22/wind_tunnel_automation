#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Windy 風洞自動化システム 課題・原因分析・今後の方針 スライド生成スクリプト

実行:
    python make_issues_pptx.py
出力:
    Windy_課題と今後.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ============================================================
#  カラーパレット（操作マニュアルと統一）
# ============================================================
C_DARK_BLUE   = RGBColor(0x1A, 0x37, 0x6C)
C_MID_BLUE    = RGBColor(0x2E, 0x6E, 0xB4)
C_LIGHT_BLUE  = RGBColor(0xD6, 0xE8, 0xF7)
C_GREEN       = RGBColor(0x21, 0x7A, 0x3C)
C_GREEN_LIGHT = RGBColor(0xD6, 0xEE, 0xDF)
C_ORANGE      = RGBColor(0xD9, 0x6D, 0x00)
C_ORANGE_LIGHT= RGBColor(0xFD, 0xF0, 0xDE)
C_GRAY_BG     = RGBColor(0xF4, 0xF6, 0xF9)
C_TEXT        = RGBColor(0x1A, 0x1A, 0x2E)
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_RED_LIGHT   = RGBColor(0xFD, 0xE8, 0xE8)
C_RED         = RGBColor(0xC0, 0x39, 0x2B)
C_SUB         = RGBColor(0x5A, 0x6A, 0x85)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ============================================================
#  ヘルパー
# ============================================================
def add_rect(slide, l, t, w, h, fill=None, border=None, border_w=Pt(1)):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.shadow.inherit = False
    if fill:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border; shape.line.width = border_w
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=Pt(14), bold=False,
             color=C_TEXT, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = size; r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return tb

def slide_header(slide, title, subtitle=None, idx=None):
    add_rect(slide, 0, Inches(0.85), prs.slide_width,
             prs.slide_height - Inches(0.85), fill=C_GRAY_BG)
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.85), fill=C_DARK_BLUE)
    add_text(slide, title, Inches(0.35), Inches(0.07), Inches(11.5), Inches(0.55),
             size=Pt(23), bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.37), Inches(0.55), Inches(11.5), Inches(0.3),
                 size=Pt(11.5), color=RGBColor(0xA8, 0xC8, 0xF0))
    if idx is not None:
        add_text(slide, idx, Inches(12.3), Inches(0.25), Inches(0.9), Inches(0.4),
                 size=Pt(13), bold=True, color=RGBColor(0xA8, 0xC8, 0xF0),
                 align=PP_ALIGN.RIGHT)

def chip(slide, text, l, t, color, w=Inches(1.15), h=Inches(0.34), tcolor=C_WHITE):
    add_rect(slide, l, t, w, h, fill=color)
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(10.5); r.font.bold = True; r.font.color.rgb = tcolor

def bullets(slide, lines, l, t, w, h, size=Pt(13), gap=Pt(5)):
    """lines: list of (text, level, color) or (text, level)"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in lines:
        text = item[0]; level = item[1] if len(item) > 1 else 0
        color = item[2] if len(item) > 2 else C_TEXT
        bold = item[3] if len(item) > 3 else False
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = gap; p.level = level
        mark = "■ " if level == 0 else "・ " if level == 1 else "– "
        r = p.add_run(); r.text = mark + text
        r.font.size = size if level == 0 else Pt(size.pt - 1.5)
        r.font.color.rgb = color; r.font.bold = bold
    return tb

def panel(slide, l, t, w, h, fill, border):
    add_rect(slide, l, t, w, h, fill=fill, border=border, border_w=Pt(1.5))


# ============================================================
#  S1 タイトル
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, prs.slide_width, prs.slide_height, fill=C_DARK_BLUE)
add_rect(s, 0, Inches(4.7), prs.slide_width, Inches(2.8), fill=RGBColor(0x14,0x2A,0x55))
add_rect(s, Inches(0.9), Inches(2.05), Inches(0.16), Inches(1.7), fill=C_MID_BLUE)
add_text(s, "風洞実験自動化システム（Windy）",
         Inches(1.2), Inches(2.0), Inches(11), Inches(0.9),
         size=Pt(36), bold=True, color=C_WHITE)
add_text(s, "現状の課題・原因分析・今後の方針",
         Inches(1.22), Inches(2.95), Inches(11), Inches(0.6),
         size=Pt(22), color=RGBColor(0xA8, 0xC8, 0xF0))
add_text(s, "Leptrino 6軸力センサ × ロータリーステージ × 差圧風速計  /  MATLAB自動計測 + Python後処理",
         Inches(1.22), Inches(5.05), Inches(11), Inches(0.5),
         size=Pt(13), color=RGBColor(0x9F,0xB6,0xD8))
add_text(s, "2026-06-08",
         Inches(1.22), Inches(5.5), Inches(6), Inches(0.4),
         size=Pt(13), color=RGBColor(0x7F,0x97,0xC0))


# ============================================================
#  S2 背景・問題の発端
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "背景：自動化後、結果が過去データと系統的にずれた", idx="1")

panel(s, Inches(0.4), Inches(1.05), Inches(6.05), Inches(2.6), C_WHITE, C_MID_BLUE)
add_text(s, "システム概要", Inches(0.55), Inches(1.15), Inches(5.7), Inches(0.4),
         size=Pt(15), bold=True, color=C_MID_BLUE)
bullets(s, [
    ("旧来の手動計測を MATLAB で自動化", 0),
    ("力: Leptrino 6軸センサ（自作 leptrino_server.py で取得）", 1),
    ("迎角: QT-ADL1 ロータリーステージ", 1),
    ("風速: 差圧センサ電圧 → 風速（make_windspeed.py）", 1),
    ("後処理: calc_force.py で空力係数 Cl, Cd, Cm を算出", 1),
], Inches(0.55), Inches(1.55), Inches(5.8), Inches(2.0), size=Pt(12.5))

panel(s, Inches(6.65), Inches(1.05), Inches(6.25), Inches(2.6), C_RED_LIGHT, C_RED)
add_text(s, "問題の発端", Inches(6.8), Inches(1.15), Inches(5.9), Inches(0.4),
         size=Pt(15), bold=True, color=C_RED)
bullets(s, [
    ("自動化システムの計測結果が、過去の手動計測と", 0),
    ("系統的に一致しない", 1, C_RED, True),
    ("同一の NACA0012 剛体翼・同一機材のはず", 0),
    ("揚力係数 Cl が過去比で大きく出る／Cm がほぼ0だった", 0),
], Inches(6.8), Inches(1.55), Inches(5.95), Inches(2.0), size=Pt(12.5))

panel(s, Inches(0.4), Inches(3.85), Inches(12.5), Inches(3.0), C_WHITE, C_MID_BLUE)
add_text(s, "調査の進め方：症状を分解 → 容疑者を1つずつ実測で潰す",
         Inches(0.55), Inches(3.95), Inches(12), Inches(0.4),
         size=Pt(15), bold=True, color=C_DARK_BLUE)
bullets(s, [
    ("「処理・較正・センサ読み込み・物理」のどこに原因があるかを、データと実測で切り分け", 0),
    ("付随して見つかった複数のバグ・不整合は順次修正（次ページ）", 0),
    ("最終的に1つの未解決の核心的問題（揚力 Fy の系統的増幅）に収束した", 0, C_RED, True),
], Inches(0.55), Inches(4.4), Inches(12.1), Inches(2.3), size=Pt(13.5), gap=Pt(8))


# ============================================================
#  S3 解決済みの問題
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "これまでに解決・対応した問題", "調査の過程で発見・修正したバグと不整合", idx="2")

rows = [
    ("モーメント単位換算バグ", "Nm→gf10cm 換算が /10（誤）→ Cm が 100倍小さく算出",
     "係数を *10 に修正。全実験フォルダの calc_force.py に適用", C_GREEN, "解決済"),
    ("plot_PM 添字ハードコード", "max_angle=20 固定で 30°計測時にグラフ破綻",
     "データ長から動的に算出するよう修正", C_GREEN, "解決済"),
    ("グラフ表示範囲", "±20° 固定で 30°まで描画されない",
     "全 calc_force.py を ±30° に統一", C_GREEN, "解決済"),
    ("ステージ原点のズレ", "ORIGIN_PULSE=11250 で過去(11025)と迎角基準が不一致",
     "11025 に戻し、過去実験と整合。NACA0012 の α₀≈0 を基準化", C_GREEN, "解決済"),
    ("ゼロ揚力角の手動管理", "取付角ドリフトの確認・補正が手作業",
     "calc_force.py に α₀ 自動レポート＋次回 ORIGIN_PULSE 推奨値を追加", C_GREEN, "解決済"),
    ("リポジトリの重複・残骸", "「 2」重複ファイル、追跡されたテスト出力",
     "重複3件を削除、pressure_data.csv を追跡解除", C_GREEN, "解決済"),
]
y = Inches(1.15)
rh = Inches(0.86)
add_rect(s, Inches(0.4), y, Inches(12.5), Inches(0.36), fill=C_DARK_BLUE)
for label, x, w in [("項目", 0.5, 2.7), ("内容（問題）", 3.3, 4.4), ("対応", 7.9, 4.0), ("状態", 11.95, 0.85)]:
    add_text(s, label, Inches(x), y+Inches(0.02), Inches(w), Inches(0.32),
             size=Pt(12), bold=True, color=C_WHITE)
y = Inches(1.55)
for i, (item, prob, fix, col, st) in enumerate(rows):
    bg = C_WHITE if i % 2 == 0 else RGBColor(0xEC,0xF1,0xF8)
    add_rect(s, Inches(0.4), y, Inches(12.5), rh, fill=bg, border=RGBColor(0xD0,0xDA,0xE8), border_w=Pt(0.75))
    add_text(s, item, Inches(0.5), y+Inches(0.06), Inches(2.7), rh-Inches(0.1),
             size=Pt(11.5), bold=True, color=C_DARK_BLUE)
    add_text(s, prob, Inches(3.3), y+Inches(0.06), Inches(4.5), rh-Inches(0.1), size=Pt(10.5))
    add_text(s, fix, Inches(7.9), y+Inches(0.06), Inches(4.0), rh-Inches(0.1), size=Pt(10.5))
    chip(s, st, Inches(11.95), y+Inches(0.25), C_GREEN, w=Inches(0.82), h=Inches(0.32))
    y += rh


# ============================================================
#  S4 未解決の主問題（症状）
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "未解決の核心的問題：揚力 Fy の系統的な増幅", "壁補正前の揚力傾斜が過去比 1.29 倍", idx="3")

panel(s, Inches(0.4), Inches(1.05), Inches(6.15), Inches(3.0), C_WHITE, C_MID_BLUE)
add_text(s, "症状：揚力傾斜 dCl/dα が新旧で系統的に違う",
         Inches(0.55), Inches(1.15), Inches(5.85), Inches(0.6),
         size=Pt(14), bold=True, color=C_DARK_BLUE)
bullets(s, [
    ("旧 4実験：3.73〜3.77 /rad に集中", 0, C_GREEN),
    ("新 2実験：4.86〜4.89 /rad に集中", 0, C_RED),
    ("比 ≈ 1.29 倍", 0, C_RED, True),
    ("各システム内では完全に再現性あり", 1),
    ("→ ランダム誤差でなく系統的な較正定数差", 1),
], Inches(0.55), Inches(1.75), Inches(5.85), Inches(2.2), size=Pt(13))

panel(s, Inches(6.75), Inches(1.05), Inches(6.15), Inches(3.0), C_ORANGE_LIGHT, C_ORANGE)
add_text(s, "なぜ「物理的におかしい」と言えるか",
         Inches(6.9), Inches(1.15), Inches(5.85), Inches(0.6),
         size=Pt(14), bold=True, color=C_ORANGE)
bullets(s, [
    ("同一 NACA0012 翼・近いレイノルズ数（16万 vs 15万）", 0),
    ("同一翼なら dCl/dα は固定値のはず", 0),
    ("Cl は無次元化済 → 風速差ではキャンセルされる", 0),
    ("なのに 1.29 倍 = どこかに較正定数の不一致", 0, C_RED, True),
], Inches(6.9), Inches(1.75), Inches(5.85), Inches(2.2), size=Pt(13))

panel(s, Inches(0.4), Inches(4.25), Inches(12.5), Inches(2.55), C_WHITE, C_RED)
add_text(s, "派生する影響（すべて Fy 増幅に起因）",
         Inches(0.55), Inches(4.35), Inches(12), Inches(0.4),
         size=Pt(14), bold=True, color=C_RED)
bullets(s, [
    ("揚力 Cl：過去比 約1.29倍（最大の症状）", 0),
    ("抗力 Cd：中〜高迎角で水増し（D=−Fx·cosα−Fy·sinα の Fy項に混入。α=15°で抗力の約44%が偽）", 0),
    ("Cm がほぼ0 だった件：別原因（モーメント単位バグ）で、これは解決済み", 0, C_GREEN),
    ("取付角 α₀ のズレ（旧0.68°→新1.45°）：別系統の問題（ステージ原点）で対応の考え方は確立", 0, C_SUB),
], Inches(0.55), Inches(4.8), Inches(12.1), Inches(1.9), size=Pt(12.5), gap=Pt(6))


# ============================================================
#  S5 診断：6軸分解
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "診断：重心補正前の生6力を分解 → Fy 一点に収束", "風速差(q)を補正した「力そのもの」のスケール（新/旧）", idx="4")

# 表
heads = ["軸", "スケール", "信頼度", "判定"]
xs    = [0.5, 1.9, 3.6, 5.3]
ws    = [1.3, 1.6, 1.6, 2.0]
data = [
    ("Fy（揚力）", "1.29", "高 r=0.999", "増幅（本体）", C_RED),
    ("Fz", "1.23", "—", "微小側力＝ノイズ", C_SUB),
    ("Mx", "1.29", "高 r=0.999", "= Fy×アーム(13cm)", C_SUB),
    ("Fx（抗力）", "1.07", "高 r=0.998", "ほぼ正常", C_GREEN),
    ("My", "1.07", "高 r=0.998", "ほぼ正常", C_GREEN),
    ("Mz", "0.86", "低 r=0.68", "ほぼ0・ノイズ", C_SUB),
]
y0 = Inches(1.15)
add_rect(s, Inches(0.4), y0, Inches(7.2), Inches(0.36), fill=C_DARK_BLUE)
for hx, hw, ht in zip(xs, ws, heads):
    add_text(s, ht, Inches(hx), y0+Inches(0.02), Inches(hw), Inches(0.32),
             size=Pt(11.5), bold=True, color=C_WHITE)
yy = Inches(1.55)
for i,(ax, sc, rel, jd, col) in enumerate(data):
    bg = C_WHITE if i%2==0 else RGBColor(0xEC,0xF1,0xF8)
    add_rect(s, Inches(0.4), yy, Inches(7.2), Inches(0.5), fill=bg,
             border=RGBColor(0xD0,0xDA,0xE8), border_w=Pt(0.75))
    add_text(s, ax, Inches(0.5), yy+Inches(0.08), Inches(1.4), Inches(0.36), size=Pt(11.5), bold=True, color=col)
    add_text(s, sc+"×", Inches(1.9), yy+Inches(0.08), Inches(1.6), Inches(0.36), size=Pt(11.5), bold=True, color=col)
    add_text(s, rel, Inches(3.6), yy+Inches(0.08), Inches(1.6), Inches(0.36), size=Pt(10.5))
    add_text(s, jd, Inches(5.3), yy+Inches(0.08), Inches(2.2), Inches(0.36), size=Pt(10.5), color=col)
    yy += Inches(0.5)

panel(s, Inches(7.85), Inches(1.15), Inches(5.05), Inches(3.5), C_RED_LIGHT, C_RED)
add_text(s, "結論：すべて『揚力 Fy 1.29倍』に集約",
         Inches(8.0), Inches(1.25), Inches(4.75), Inches(0.6),
         size=Pt(14), bold=True, color=C_RED)
bullets(s, [
    ("Mx は Fy×アーム（新旧でアーム長 0.130m 一致）→ Fy の継承", 0),
    ("Fz・Mz はほぼ0の微小量でノイズ（倍率に意味なし）", 0),
    ("Fx・My（抗力方向）は正常", 0, C_GREEN),
    ("抗力にも Fy が sinα で混入（自己整合）", 0),
    ("→ 独立した異常は Fy ただ一つ", 0, C_RED, True),
], Inches(8.0), Inches(1.85), Inches(4.75), Inches(2.7), size=Pt(12.5), gap=Pt(7))

panel(s, Inches(0.4), Inches(4.7), Inches(7.2), Inches(2.1), C_WHITE, C_MID_BLUE)
add_text(s, "重要：揚力(Fy)は水平方向、重量(Fz)は鉛直方向",
         Inches(0.55), Inches(4.8), Inches(6.9), Inches(0.4), size=Pt(13), bold=True, color=C_DARK_BLUE)
bullets(s, [
    ("翼は鉛直軸まわりに回転 → 重量26NはFz軸（鉛直）", 0),
    ("揚力は水平面内 → Fy軸", 0),
    ("おもり(鉛直)で検証したのは Fz軸。Fy軸は未検証", 0, C_RED, True),
], Inches(0.55), Inches(5.2), Inches(6.9), Inches(1.5), size=Pt(12))


# ============================================================
#  S6 切り分け：潰した容疑者
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "切り分け：処理・較正・読み込みは『シロ』", "実測・コード照合で1つずつ除外", idx="5")

checks = [
    ("力センサ（静荷重）", "既知おもりを 0.25% 精度で正しく計測", "Fz軸は正確"),
    ("差圧電圧 mV", "自動取得値が手動読みと一致", "読み取りOK"),
    ("風速計算式", "旧Excel と make_windspeed.py が小数8桁まで一致", "式は同一"),
    ("気温・気圧 → 空気密度ρ", "ρは q=0.5ρU²S で打ち消され Cl に効かない（実証）", "無関係"),
    ("力→係数 後処理", "calc_force.py の較正行列・換算が新旧バイト一致", "処理同一"),
    ("迎角の正しさ", "失速角が新旧で整合（角度スケール誤差なし）", "角度OK"),
    ("揚抗力の三角関数", "壁補正前の幾何迎角を使用（正しい）。生Fyより下流", "問題なし"),
    ("データロガー変換", "定格/10000 換算・軸順・ヘッダが新旧同一", "差分なし"),
]
col_w = Inches(6.1)
for i,(item, detail, verdict) in enumerate(checks):
    col = i % 2
    row = i // 2
    l = Inches(0.4) + col * (col_w + Inches(0.3))
    t = Inches(1.15) + row * Inches(1.32)
    panel(s, l, t, col_w, Inches(1.18), C_GREEN_LIGHT, C_GREEN)
    add_text(s, "✓ " + item, l+Inches(0.12), t+Inches(0.08), col_w-Inches(1.4), Inches(0.4),
             size=Pt(13), bold=True, color=C_GREEN)
    chip(s, verdict, l+col_w-Inches(1.25), t+Inches(0.1), C_GREEN, w=Inches(1.15), h=Inches(0.3))
    add_text(s, detail, l+Inches(0.12), t+Inches(0.5), col_w-Inches(0.25), Inches(0.62), size=Pt(11.5))

add_text(s, "→ 自動化の数値処理・センサ読み込み・較正に、Fy を選択的にずらすバグは発見されず",
         Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.5),
         size=Pt(14), bold=True, color=C_DARK_BLUE, align=PP_ALIGN.CENTER)


# ============================================================
#  S7 残る容疑と解決法の模索
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "残る容疑と解決法の模索", "Fy だけが・高プリロード下で・乗法的に増幅", idx="6")

panel(s, Inches(0.4), Inches(1.05), Inches(6.15), Inches(3.55), C_ORANGE_LIGHT, C_ORANGE)
add_text(s, "容疑① 軸間干渉（最有力）",
         Inches(0.55), Inches(1.15), Inches(5.85), Inches(0.4), size=Pt(15), bold=True, color=C_ORANGE)
bullets(s, [
    ("旧ロガー：計測前にタレ（ゼロ点）→ センサはゼロ付近で動作", 0, C_GREEN),
    ("新システム：タレ無し → Fz=−26N（定格30Nの87%）で常時動作", 0, C_RED),
    ("大きな Fz プリロードが Fy のゲインを変化させる", 0, C_RED, True),
    ("（6軸センサの combined-loading interference）", 1, C_SUB),
    ("ゲイン変化は乗法的 → Pofst 減算で消えない", 0),
    ("観測（Fy増幅・Fx正常・Fzで検証済）と機構的に整合", 0),
], Inches(0.55), Inches(1.6), Inches(5.85), Inches(2.9), size=Pt(12), gap=Pt(6))

panel(s, Inches(6.75), Inches(1.05), Inches(6.15), Inches(3.55), C_WHITE, C_MID_BLUE)
add_text(s, "容疑② 物理マウントのたわみ",
         Inches(6.9), Inches(1.15), Inches(5.85), Inches(0.4), size=Pt(15), bold=True, color=C_MID_BLUE)
bullets(s, [
    ("揚力（横荷重）下で支持系・スティングがたわむ", 0),
    ("実効迎角が荷重に応じて増える → 揚力方向に集中", 0),
    ("剛体翼でも支持系の剛性差なら起こり得る", 0),
    ("新リグの結合・剛性が旧と異なる可能性", 0),
], Inches(6.9), Inches(1.6), Inches(5.85), Inches(2.0), size=Pt(12.5))
add_text(s, "補足：ソフト的タレは結果を変えない（Pofst減算で相殺）→ 物理対策が必要",
         Inches(6.9), Inches(3.75), Inches(5.85), Inches(0.8), size=Pt(11.5), italic=True, color=C_SUB)

panel(s, Inches(0.4), Inches(4.8), Inches(12.5), Inches(2.0), C_LIGHT_BLUE, C_MID_BLUE)
add_text(s, "決定的な切り分け実験：Fy 軸 水平荷重テスト",
         Inches(0.55), Inches(4.9), Inches(12), Inches(0.4), size=Pt(14), bold=True, color=C_DARK_BLUE)
bullets(s, [
    ("①翼装着(Fz≈−26N)で Fy 方向に既知荷重 → 1.29倍出れば「Fy軸ゲイン/干渉」確定", 0),
    ("②翼を外し Fz≈0 で同荷重 → 正確なら「Fzプリロードによる干渉」が確定", 0),
    ("ツール: tare_measure.m（ゼロ点取得→相対6力表示）を実装済み", 0, C_GREEN, True),
], Inches(0.55), Inches(5.32), Inches(12.1), Inches(1.4), size=Pt(12.5), gap=Pt(6))


# ============================================================
#  S8 今後やること
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "今後やること（ロードマップ）", idx="7")

steps = [
    ("STEP 1", "実験室で Fy 軸 水平荷重テスト", "プリロード有/無で実施。tare_measure.m を使用。最優先", C_RED, "次回実験"),
    ("STEP 2", "原因の確定", "①Fy軸ゲイン/軸間干渉 か ②物理マウント か を判定", C_ORANGE, "テスト後"),
    ("STEP 3", "対策の実施", "干渉→カウンターウェイトで Fz プリロード低減／センサ再較正。たわみ→支持系の剛性向上", C_MID_BLUE, "確定後"),
    ("STEP 4", "（補助）ピトー管で風速の独立検証", "差圧タップ由来 U と直接比較し、動圧側の妥当性も確認", C_SUB, "並行"),
    ("STEP 5", "全実験データの再処理・グラフ更新", "原因解消後、過去・新データを統一条件で再算出し検証", C_GREEN, "最終"),
]
y = Inches(1.2)
for tag, title, detail, col, when in steps:
    add_rect(s, Inches(0.4), y, Inches(1.55), Inches(1.0), fill=col)
    tb = s.shapes.add_textbox(Inches(0.4), y, Inches(1.55), Inches(1.0))
    tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = tag; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = C_WHITE
    add_rect(s, Inches(2.05), y, Inches(10.85), Inches(1.0), fill=C_WHITE,
             border=RGBColor(0xD0,0xDA,0xE8), border_w=Pt(1))
    add_text(s, title, Inches(2.25), y+Inches(0.12), Inches(8.0), Inches(0.45),
             size=Pt(15), bold=True, color=col)
    add_text(s, detail, Inches(2.25), y+Inches(0.55), Inches(9.0), Inches(0.4),
             size=Pt(11.5), color=C_TEXT)
    chip(s, when, Inches(11.55), y+Inches(0.33), col, w=Inches(1.2), h=Inches(0.34))
    y += Inches(1.12)


# ============================================================
#  S9 まとめ
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, prs.slide_width, prs.slide_height, fill=C_DARK_BLUE)
add_rect(s, Inches(0.9), Inches(0.85), Inches(0.16), Inches(0.7), fill=C_MID_BLUE)
add_text(s, "まとめ", Inches(1.2), Inches(0.8), Inches(10), Inches(0.8),
         size=Pt(30), bold=True, color=C_WHITE)

panel(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(1.45), RGBColor(0x14,0x2A,0x55), C_MID_BLUE)
bullets(s, [
    ("付随する複数のバグ・不整合は修正済み（モーメント単位、グラフ範囲、ステージ原点、α₀レポート、整理）", 0, C_WHITE),
    ("Cm≈0 や Cl の一部は解決。残る核心は『揚力 Fy が過去比 1.29倍』ただ一点", 0, RGBColor(0xFF,0xD0,0xC0), True),
], Inches(1.15), Inches(2.05), Inches(11.0), Inches(1.2), size=Pt(13.5), gap=Pt(8))

panel(s, Inches(0.9), Inches(3.55), Inches(11.5), Inches(1.45), RGBColor(0x14,0x2A,0x55), C_GREEN)
bullets(s, [
    ("処理・較正・読み込みは網羅的に検証し『シロ』を確認", 0, C_WHITE),
    ("残る容疑は ①軸間干渉（タレ無し・Fz高プリロード）②物理マウントのたわみ の2系統", 0, RGBColor(0xC8,0xF0,0xD8)),
], Inches(1.15), Inches(3.7), Inches(11.0), Inches(1.2), size=Pt(13.5), gap=Pt(8))

panel(s, Inches(0.9), Inches(5.2), Inches(11.5), Inches(1.5), RGBColor(0x2A,0x1A,0x10), C_ORANGE)
add_text(s, "次の一手", Inches(1.15), Inches(5.3), Inches(10), Inches(0.4),
         size=Pt(15), bold=True, color=C_ORANGE)
bullets(s, [
    ("実験室で Fy 軸 水平荷重テスト（プリロード有/無）→ 原因を確定 → 物理対策へ", 0, C_WHITE, True),
    ("ツール tare_measure.m 実装済み。確定後に全データ再処理", 0, RGBColor(0xF5,0xD8,0xB0)),
], Inches(1.15), Inches(5.7), Inches(11.0), Inches(1.0), size=Pt(13.5), gap=Pt(7))


# ============================================================
prs.save("Windy_課題と今後.pptx")
print("生成完了: Windy_課題と今後.pptx  (全%dスライド)" % len(prs.slides._sldIdLst))
