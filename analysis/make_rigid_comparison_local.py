#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
make_rigid_comparison_pptx.py
NACA0012 rigid wing 全データ比較 パワポ生成
Ito(240521/240603/241223) + 250924/251020 + 260417/260424/260430/260520 + 260605/260605_2
260605/260605_2 は風洞自動化新システムによる計測
"""

import os, io, struct, sys, json
import numpy as np
import pandas as pd

# 端末/MATLAB の system() 経由でも文字エンコードで落ちないようにする安全網。
for _stream in (sys.stdout, sys.stderr):
    try:
        _stream.reconfigure(errors="backslashreplace")
    except (AttributeError, ValueError):
        pass
from scipy.interpolate import interp1d
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.font_manager as _fm
# 利用可能な日本語フォントを自動選択（Mac=Hiragino / Windows=Yu Gothic 等）
_avail = {f.name for f in _fm.fontManager.ttflist}
for _jp in ["Hiragino Sans", "Yu Gothic", "Meiryo", "MS Gothic",
            "Noto Sans CJK JP", "IPAexGothic", "TakaoGothic"]:
    if _jp in _avail:
        plt.rcParams["font.family"] = _jp
        break
import matplotlib.ticker as ticker
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn as _qn

# ─── 定数 ────────────────────────────────────────────────────────────────────
SCRIPT_DIR    = os.path.dirname(os.path.abspath(__file__))
BASE          = os.path.join(SCRIPT_DIR, "aero_data")   # 各実験 C_aero.csv の置き場
TEMPLATE_PATH = os.path.join(SCRIPT_DIR, "研究室MTGテンプレート.pptx")

def _output_dir():
    """config.json の output_dir を返す。無ければ SCRIPT_DIR にフォールバック。"""
    config_path = os.path.join(os.path.dirname(SCRIPT_DIR), "config.json")
    try:
        with open(config_path, encoding="utf-8") as f:
            d = json.load(f)
        out = d.get("output_dir", "")
        if out:
            os.makedirs(out, exist_ok=True)
            return out
    except (OSError, json.JSONDecodeError):
        pass
    return SCRIPT_DIR

OUT = os.path.join(_output_dir(), "Windy新システムによる実験結果.pptx")

SLIDE_W = Inches(10.0)
SLIDE_H = Inches(7.5)

# ─── データ定義 ───────────────────────────────────────────────────────────────
# 色系統（3カテゴリ）：旧システム=寒色（青系）／新システム接触不良=暖色（赤橙系）／接触不良修正後=緑系
DATASETS = {
    "Ito\n240521":  {"path": os.path.join(BASE, "force_measurement_240521_ito_rigid/C_aero.csv"),
                     "color": "#1B3A6B", "marker": "v",  "date": "2024-05-21", "who": "old"},
    "Ito\n240603":  {"path": os.path.join(BASE, "force_measurement_240603_ito_rigid/C_aero.csv"),
                     "color": "#234E8C", "marker": "<",  "date": "2024-06-03", "who": "old"},
    "Ito\n241223":  {"path": os.path.join(BASE, "force_measurement_241223_ito_rigid/C_aero.csv"),
                     "color": "#2E63AD", "marker": ">",  "date": "2024-12-23", "who": "old"},
    "250924":       {"path": os.path.join(BASE, "force_measurement_250924_rigid/C_aero.csv"),
                     "color": "#3D7AC4", "marker": "^",  "date": "2025-09-24", "who": "old"},
    "251020":       {"path": os.path.join(BASE, "force_measurement_251020_rigid/C_aero.csv"),
                     "color": "#5591D1", "marker": "o",  "date": "2025-10-20", "who": "old"},
    "260424":       {"path": os.path.join(BASE, "force_measurement_260424/C_aero.csv"),
                     "color": "#6FA3D8", "marker": "D",  "date": "2026-04-24", "who": "old"},
    "260430":       {"path": os.path.join(BASE, "force_measurement_260430_rigid/C_aero.csv"),
                     "color": "#4472B0", "marker": "P",  "date": "2026-04-30", "who": "old"},
    "260520":       {"path": os.path.join(BASE, "force_measurement_260520_rigid/C_aero.csv"),
                     "color": "#88B4DE", "marker": "h",  "date": "2026-05-20", "who": "old"},
    "260605①":     {"path": os.path.join(BASE, "force_measurement_260605_rigid/C_aero.csv"),
                     "color": "#B03A2E", "marker": "X",  "date": "2026-06-05", "who": "new_fault"},
    "260605②":     {"path": os.path.join(BASE, "force_measurement_260605_rigid2/C_aero.csv"),
                     "color": "#E74C3C", "marker": "P",  "date": "2026-06-05", "who": "new_fault"},
    "260608":       {"path": os.path.join(BASE, "force_measurement_260608_rigid/C_aero.csv"),
                     "color": "#E67E22", "marker": "p",  "date": "2026-06-08", "who": "new_fault"},
    "260610①":     {"path": os.path.join(BASE, "force_measurement_260610_rigid/C_aero.csv"),
                     "color": "#F39C12", "marker": "d",  "date": "2026-06-10", "who": "new_fault"},
    "rigid3\n修正後": {"path": os.path.join(BASE, "force_measurement_260610_rigid3/C_aero.csv"),
                     "color": "#1E8449", "marker": "*",  "date": "2026-06-10", "who": "fixed"},
}

DISP_NAMES = {
    "Ito\n240521":  "240521",
    "Ito\n240603":  "240603",
    "Ito\n241223":  "241223",
    "250924":       "250924",
    "251020":       "251020",
    "260424":       "260424",
    "260430":       "260430",
    "260520":       "260520",
    "260605①":     "260605_1",
    "260605②":     "260605_2",
    "260608":       "260608",
    "260610①":     "260610_1",
    "rigid3\n修正後": "260610_2",
}

ROW_INFO = [
    ("240521", "old"),
    ("240603", "old"),
    ("241223", "old"),
    ("250924", "old"),
    ("251020", "old"),
    ("260424", "old"),
    ("260430", "old"),
    ("260520", "old"),
    ("260605_1", "new_fault"),
    ("260605_2", "new_fault"),
    ("260608", "new_fault"),
    ("260610_1", "new_fault"),
    ("260610_2", "fixed"),
]

# カテゴリ別 背景色・備考（表の色分けもこの3系統に統一）
BG_COLOR = {
    "old":       RGBColor(0xE3, 0xEC, 0xFB),   # 淡い青
    "new_fault": RGBColor(0xFC, 0xE6, 0xDD),   # 淡い赤橙
    "fixed":     RGBColor(0xD6, 0xF5, 0xD6),   # 淡い緑
}
NOTE_STR = {
    "old":       "（旧システム）",
    "new_fault": "（新システム・接触不良）",
    "fixed":     "（接触不良修正後）",
}

# 除外データ（aero_data には残すが図・表には載せない）
EXCLUDE_DIRS = {
    "force_measurement_260417",
    "force_meausrement_260610_rigid",   # typo名（修正済み 260610_rigid と同内容の再同期を防ぐ）
}

# ─── aero_data 内の未登録フォルダを自動追加 ──────────────────────────────────────
#  今後 update_aero_data.py で aero_data/ に新しい rigid 実験の C_aero.csv が
#  追加されると、ここで自動的に DATASETS へ取り込まれ、表・全図に反映される。
#  （既存の登録済み13件はそのままのスタイル・解説を維持）
#  自動追加データは who="new"（新システム・新規）として扱い、緑〜青緑系の色を割当てる。
#  who="new" は揚力傾斜−過去平均スライドの new_keys にも含めるため、表・全図に加え
#  「過去平均との差」グラフにも自動で反映される（過去平均の母集団は who=="old" のみ）。
import re as _re
_known_paths = {os.path.normpath(v["path"]) for v in DATASETS.values()}
_auto_palette = ["#16A085", "#27AE60", "#2ECC71", "#1ABC9C", "#138D75",
                 "#52BE80", "#45B39D", "#229954", "#117A65", "#0E6655"]
_auto_markers = ["o", "s", "^", "D", "v", "p", "X", "h", "<", ">"]
_ai = 0
if os.path.isdir(BASE):
    for _sub in sorted(os.listdir(BASE)):
        if _sub in EXCLUDE_DIRS:
            continue
        _ca = os.path.join(BASE, _sub, "C_aero.csv")
        if not os.path.isfile(_ca):
            continue
        if os.path.normpath(_ca) in _known_paths:
            continue
        _m    = _re.search(r"(\d{6})", _sub)
        _date = f"20{_m.group(1)[:2]}-{_m.group(1)[2:4]}-{_m.group(1)[4:]}" if _m else "—"
        DATASETS[_sub]   = {"path": _ca,
                            "color":  _auto_palette[_ai % len(_auto_palette)],
                            "marker": _auto_markers[_ai % len(_auto_markers)],
                            "date":   _date, "who": "new"}
        # フォルダ名 → 表示名（数字のみ。typo 綴り force_meausrement_ も吸収）
        _disp = (_sub.replace("force_measurement_", "").replace("force_meausrement_", "")
                     .replace("_rigid", ""))
        DISP_NAMES[_sub] = _disp
        ROW_INFO.append((_disp, "new"))
        _ai += 1
BG_COLOR.setdefault("new", RGBColor(0xDD, 0xF5, 0xE8))   # 淡い緑（新システム新規）
NOTE_STR.setdefault("new", "（新システム・新規）")
if _ai:
    print(f"[自動追加] aero_data 内の未登録 {_ai} 件を取り込みました。")

# ─── データ読み込み ───────────────────────────────────────────────────────────
def load(info):
    df = pd.read_csv(info["path"])
    df = df.apply(pd.to_numeric, errors="coerce").dropna()
    return df

data = {k: load(v) for k, v in DATASETS.items()}

# ─── 統計量計算 ───────────────────────────────────────────────────────────────
stats = {}
for k, df in data.items():
    row0  = df[df["AoA"] == 0]
    cl0   = float(row0.iloc[0]["Cl"]) if len(row0) > 0 else np.nan
    cd0   = float(row0.iloc[0]["Cd"]) if len(row0) > 0 else np.nan
    cm0   = float(row0.iloc[0]["Cm"]) if len(row0) > 0 else np.nan
    seg   = df[(df["AoA"] >= -5) & (df["AoA"] <= 5)].sort_values("AoA_mod")
    slope = float(np.polyfit(seg["AoA_mod"].astype(float),
                              seg["Cl"].astype(float), 1)[0]) if len(seg) >= 2 else np.nan
    # 零揚力迎角
    aoa0 = None
    for i in range(len(seg) - 1):
        c1, c2 = float(seg.iloc[i]["Cl"]), float(seg.iloc[i+1]["Cl"])
        if c1 <= 0 <= c2 or c2 <= 0 <= c1:
            am1 = float(seg.iloc[i]["AoA_mod"])
            am2 = float(seg.iloc[i+1]["AoA_mod"])
            aoa0 = am1 + (0 - c1) / (c2 - c1) * (am2 - am1)
            break
    if aoa0 is None:
        aoa0 = 0.0
    # P-M 対称性
    pos = df[df["AoA"] >= 0].sort_values("AoA")
    neg = df[df["AoA"] <= 0].sort_values("AoA")
    if len(neg) > 3:
        aoa_pm = np.arange(1, 21, dtype=float)
        fp = np.interp(aoa_pm, pos["AoA"].values, pos["Cl"].values,
                       left=np.nan, right=np.nan)
        fn = np.interp(aoa_pm, -neg["AoA"].values[::-1], neg["Cl"].values[::-1],
                       left=np.nan, right=np.nan)
        diff = fp - (-fn)
        mask = ~np.isnan(diff)
        pm_rmse = float(np.sqrt(np.mean(diff[mask]**2))) if mask.any() else np.nan
    else:
        pm_rmse = np.nan
    stats[k] = {"Cl0": cl0, "Cd0": cd0, "Cm0": cm0,
                 "aoa0": aoa0, "slope": slope, "pm_rmse": pm_rmse,
                 "date": DATASETS[k]["date"], "color": DATASETS[k]["color"]}

# ─── 零揚力迎角補正後 P-M RMSE ───────────────────────────────────────────────
pm_rmse_corr = {}
for k, df in data.items():
    df_s = df.sort_values("AoA_mod")
    a0 = stats[k]["aoa0"]
    deltas = np.arange(1, 21, dtype=float)
    f = interp1d(df_s["AoA_mod"].astype(float), df_s["Cl"].astype(float),
                 kind="linear", bounds_error=False, fill_value=np.nan)
    v_plus  = f(a0 + deltas)
    v_minus = f(a0 - deltas)
    diff = v_plus - (-v_minus)
    mask = ~(np.isnan(v_plus) | np.isnan(v_minus))
    pm_rmse_corr[k] = float(np.sqrt(np.mean(diff[mask]**2))) if mask.any() else np.nan

labels  = list(DATASETS.keys())
x_pos   = np.arange(len(labels))
colors  = [DATASETS[k]["color"] for k in labels]
disp    = [DISP_NAMES.get(k, k.replace("\n", " ")) for k in labels]

# ─── 図作成ユーティリティ ─────────────────────────────────────────────────────
def fig_to_bytes(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=180, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf

TS  = 12
LW  = 1.5
MK  = 4

# 特定データの強調はしない（全データを同等に表示）
def _is_emph(k):
    return False

# ─── 図1: 零揚力迎角の時系列 ─────────────────────────────────────────────────
fig, ax = plt.subplots(figsize=(10, 5))
vals = [stats[k]["aoa0"] for k in labels]
bar_colors = [DATASETS[k]["color"] for k in labels]
bars = ax.bar(x_pos, vals, color=bar_colors, edgecolor="white", linewidth=0.8, width=0.7)
for bar, v in zip(bars, vals):
    ax.text(bar.get_x() + bar.get_width()/2, v + 0.01,
            f"{v:.3f}°", ha="center", va="bottom", fontsize=9, fontweight="bold")
ax.axhline(0, color="black", linewidth=0.8, linestyle="--")
# 修正後（旧システム一致）を背景強調
for i, k in enumerate(labels):
    if _is_emph(k):
        ax.axvspan(i - 0.45, i + 0.45, alpha=0.15, color="mediumseagreen", zorder=0)
ax.set_xticks(x_pos)
ax.set_xticklabels(disp, fontsize=TS - 2, rotation=15, ha="right")
ax.set_ylabel("零揚力迎角 [deg]", fontsize=TS)
_vmax = max([v for v in vals if not np.isnan(v)] + [0.0])
ax.set_ylim(-0.1, _vmax * 1.18 + 0.08)   # 値ラベル分の余白を確保（見切れ防止）
ax.set_title("零揚力迎角の時系列変化（Cl=0 となる AoA_mod）", fontsize=TS)
ax.yaxis.set_major_locator(ticker.MultipleLocator(0.2))
ax.grid(axis="y", linestyle="--", alpha=0.5)
fig.tight_layout()
img_aoa0 = fig_to_bytes(fig)

# ─── 図2: Clスロープの時系列 ─────────────────────────────────────────────────
fig, ax = plt.subplots(figsize=(10, 5))
slopes = [stats[k]["slope"] for k in labels]
bars = ax.bar(x_pos, slopes, color=bar_colors, edgecolor="white", linewidth=0.8, width=0.7)
for bar, v in zip(bars, slopes):
    if not np.isnan(v):
        ax.text(bar.get_x() + bar.get_width()/2, v + 0.0005,
                f"{v:.4f}", ha="center", va="bottom", fontsize=9, fontweight="bold")
# 2π/rad = ~0.10966/deg 理論線
ax.axhline(2 * np.pi / 180, color="gray", linewidth=1.2, linestyle=":",
           label=r"$2\pi$ /rad = {:.4f}/deg".format(2*np.pi/180))
for i, k in enumerate(labels):
    if _is_emph(k):
        ax.axvspan(i - 0.45, i + 0.45, alpha=0.15, color="mediumseagreen", zorder=0)
ax.set_xticks(x_pos)
ax.set_xticklabels(disp, fontsize=TS - 2, rotation=15, ha="right")
ax.set_ylabel("dCl/dAoA_mod [/deg]", fontsize=TS)
ax.set_ylim(0.09, 0.18)
ax.set_title("線形域 Cl スロープの時系列変化（AoA_mod = −5〜+5°、線形回帰）", fontsize=TS)
ax.yaxis.set_major_locator(ticker.MultipleLocator(0.01))
ax.legend(fontsize=10)
ax.grid(axis="y", linestyle="--", alpha=0.5)
fig.tight_layout()
img_slope = fig_to_bytes(fig)

# ─── 図2b: 新システムの揚力傾斜 − 過去剛体翼(旧システム)の平均 ─────────────────
old_keys   = [k for k in labels if DATASETS[k]["who"] == "old"]
new_keys    = [k for k in labels if DATASETS[k]["who"] in ("new_fault", "fixed", "new")]
old_slopes = [stats[k]["slope"] for k in old_keys if not np.isnan(stats[k]["slope"])]
OLD_AVG    = float(np.mean(old_slopes))
OLD_STD    = float(np.std(old_slopes))
slope_diffs = [stats[k]["slope"] - OLD_AVG for k in new_keys]
slope_pcts  = [(stats[k]["slope"] / OLD_AVG - 1.0) * 100.0 for k in new_keys]

fig, ax = plt.subplots(figsize=(9, 5))
bc   = [DATASETS[k]["color"] for k in new_keys]
bars = ax.bar(range(len(new_keys)), slope_diffs, color=bc,
              edgecolor="white", linewidth=0.8, width=0.6)
# 修正後（旧と一致）を背景強調
for i, k in enumerate(new_keys):
    if _is_emph(k):
        ax.axvspan(i - 0.4, i + 0.4, alpha=0.18, color="mediumseagreen", zorder=0)
# 過去平均=0 の基準線＋過去ばらつき帯
ax.axhline(0, color="black", linewidth=1.0)
ax.axhspan(-OLD_STD, OLD_STD, alpha=0.10, color="royalblue", zorder=0,
           label=f"過去平均 {OLD_AVG:.4f}/deg ±1σ")
for bar, d, p in zip(bars, slope_diffs, slope_pcts):
    va  = "bottom" if d >= 0 else "top"
    off = 0.0010 if d >= 0 else -0.0010
    ax.text(bar.get_x() + bar.get_width()/2, d + off,
            f"{d:+.4f}\n({p:+.1f}%)", ha="center", va=va,
            fontsize=9, fontweight="bold")
ax.set_xticks(range(len(new_keys)))
_nk = len(new_keys)
ax.set_xticklabels([DISP_NAMES.get(k, k) for k in new_keys],
                   fontsize=(TS - 1 if _nk <= 7 else TS - 3),
                   rotation=(0 if _nk <= 7 else 20),
                   ha=("center" if _nk <= 7 else "right"))
ax.set_ylabel("揚力傾斜の差  Δ(dCl/dα) [/deg]", fontsize=TS)
ax.set_title(f"新システムの揚力傾斜 − 過去剛体翼の平均（{OLD_AVG:.4f}/deg）", fontsize=TS)
ax.legend(fontsize=10, loc="upper right")
ax.grid(axis="y", linestyle="--", alpha=0.5)
fig.tight_layout()
img_slope_diff = fig_to_bytes(fig)

# ─── 図3: Cl 曲線オーバーレイ ─────────────────────────────────────────────────
fig, ax = plt.subplots(figsize=(7.5, 6))
for k, info in DATASETS.items():
    df = data[k]
    lw = 2.5 if _is_emph(k) else LW
    ms = 5   if _is_emph(k) else MK
    ax.plot(df["AoA_mod"], df["Cl"],
            color=info["color"], marker=info["marker"],
            linewidth=lw, markersize=ms,
            label=DISP_NAMES.get(k, k.replace("\n", " ")))
ax.axhline(0, color="black", linewidth=0.6, linestyle="--")
ax.axvline(0, color="black", linewidth=0.6, linestyle="--")
aoa_th = np.linspace(-25, 25, 200)
ax.plot(aoa_th, 2*np.pi*np.deg2rad(aoa_th), color="gray",
        linestyle=":", linewidth=1.5, label=r"$2\pi$ 理論")
ax.set_xlim(-25, 25)
ax.set_ylim(-1.0, 1.0)
ax.set_xlabel("AoA (補正後) [deg]", fontsize=TS)
ax.set_ylabel(r"$C_l$", fontsize=TS)
ax.set_title("Cl 曲線 全データ比較", fontsize=TS)
ax.legend(fontsize=8, ncol=(2 if len(DATASETS) <= 16 else 3), loc="upper left")
ax.grid(linestyle="--", alpha=0.4)
fig.tight_layout()
img_cl = fig_to_bytes(fig)

# ─── 図4: AoA=0 付近の拡大 ───────────────────────────────────────────────────
fig, ax = plt.subplots(figsize=(7.5, 6))
for k, info in DATASETS.items():
    df = data[k]
    near = df[(df["AoA"] >= -5) & (df["AoA"] <= 5)].sort_values("AoA")
    lw = 2.5 if _is_emph(k) else LW
    ms = 6   if _is_emph(k) else 5
    ax.plot(near["AoA_mod"], near["Cl"],
            color=info["color"], marker=info["marker"],
            linewidth=lw, markersize=ms,
            label=DISP_NAMES.get(k, k.replace("\n", " ")))
ax.axhline(0, color="black", linewidth=1, linestyle="--")
ax.axvline(0, color="black", linewidth=0.6, linestyle="--")
ax.set_xlim(-5, 5)
ax.set_ylim(-0.5, 0.5)
ax.set_xlabel("AoA (補正後) [deg]", fontsize=TS)
ax.set_ylabel(r"$C_l$", fontsize=TS)
ax.set_title("AoA = −5〜+5° 拡大（零揚力迎角のずれ確認）", fontsize=TS)
ax.legend(fontsize=8, ncol=2)
ax.grid(linestyle="--", alpha=0.4)
fig.tight_layout()
img_zoom = fig_to_bytes(fig)

# ─── 図6a/b: P-M対称性誤差 ───────────────────────────────────────────────────
pm_keys_raw  = [k for k in labels if not np.isnan(stats[k]["pm_rmse"])]
pm_keys_corr = [k for k in labels if not np.isnan(pm_rmse_corr[k])]

def bar_pm(keys, vals, title):
    fig, ax = plt.subplots(figsize=(6.5, 5.0))   # 箱(ar≈1.3)に合わせ歪みなく充填
    bc = [DATASETS[k]["color"] for k in keys]
    dp = [DISP_NAMES.get(k, k.replace("\n", " ")) for k in keys]
    bars = ax.bar(range(len(keys)), vals, color=bc, edgecolor="white", linewidth=0.8, width=0.6)
    for i, k in enumerate(keys):
        if _is_emph(k):
            ax.axvspan(i - 0.4, i + 0.4, alpha=0.18, color="mediumseagreen", zorder=0)
    for bar, v in zip(bars, vals):
        ax.text(bar.get_x() + bar.get_width()/2, v + 0.001,
                f"{v:.3f}", ha="center", va="bottom", fontsize=9, fontweight="bold")
    ax.set_xticks(range(len(keys)))
    ax.set_xticklabels(dp, fontsize=TS - 3, rotation=15, ha="right")
    ax.set_ylabel("Cl P-M RMSE", fontsize=TS)
    ax.set_title(title, fontsize=TS)
    ax.grid(axis="y", linestyle="--", alpha=0.5)
    fig.tight_layout()
    return fig_to_bytes(fig)

img_pm      = bar_pm(pm_keys_raw,
                     [stats[k]["pm_rmse"] for k in pm_keys_raw],
                     "（a）幾何迎角基準")
img_pm_corr = bar_pm(pm_keys_corr,
                     [pm_rmse_corr[k] for k in pm_keys_corr],
                     "（b）零揚力迎角基準（補正後）")

# ─── 図7: 過去剛体翼の 平均・min/max帯 と 修正後新システム(260610_2) の比較 ────
#   α₀（取付角・原点パルス）の差は別問題なので、零揚力角基準（AoA−α₀）で揃えて
#   純粋な空力特性（傾き・形状）を比較する。
AOA_ENV = np.linspace(-20, 20, 81)         # 零揚力角基準の共通グリッド
_old_stack = []
for k in old_keys:
    df  = data[k].sort_values("AoA_mod")
    x0  = df["AoA_mod"].astype(float) - stats[k]["aoa0"]   # α₀補正
    _old_stack.append(np.interp(AOA_ENV, x0, df["Cl"].astype(float),
                                left=np.nan, right=np.nan))
_old_arr = np.vstack(_old_stack)
_cnt     = np.sum(~np.isnan(_old_arr), axis=0)
_valid   = _cnt >= max(2, len(old_keys) // 2)   # 過半数の過去データが存在する範囲のみ
with np.errstate(invalid="ignore"):
    env_mean = np.nanmean(_old_arr, axis=0)
    env_min  = np.nanmin(_old_arr, axis=0)
    env_max  = np.nanmax(_old_arr, axis=0)
xe, ymin, ymax, ymean = AOA_ENV[_valid], env_min[_valid], env_max[_valid], env_mean[_valid]

# 過去平均を 0 とした偏差（残差）で表示する
r_min = ymin - ymean
r_max = ymax - ymean

# 重ねる新システムデータ＝接触不良修正後(fixed)＋今後の新規(new)。
#   接触不良(new_fault)は既知の異常なので帯比較からは除外する。
#   → update_aero_data.py で aero_data に新実験が増えると who="new" で自動追加され、
#     このスライドにも自動で重なる。
overlay_keys = [k for k in labels if DATASETS[k]["who"] in ("fixed", "new")]

fig, ax = plt.subplots(figsize=(8, 6))
ax.fill_between(xe, r_min, r_max, color="#AFC6E8", alpha=0.6,
                label="過去剛体翼 min–max（平均からの偏差）")
ax.axhline(0, color="#234E8C", linewidth=2.0, label="過去剛体翼 平均（基準=0）")
_yvals = [r_min, r_max]
for k in overlay_keys:
    dff  = data[k].sort_values("AoA_mod")
    xf   = dff["AoA_mod"].astype(float) - stats[k]["aoa0"]
    r_f  = np.interp(xe, xf, dff["Cl"].astype(float), left=np.nan, right=np.nan) - ymean
    _yvals.append(r_f[~np.isnan(r_f)])
    ax.plot(xe, r_f, color=DATASETS[k]["color"], marker=DATASETS[k]["marker"],
            markersize=7, linewidth=2.0, label=DISP_NAMES.get(k, k))
ax.axvline(0, color="black", linewidth=0.6, linestyle="--")
# y軸はデータに合わせて自動調整
_concat = np.concatenate([v for v in _yvals if len(v)])
_ymax   = float(np.nanmax(np.abs(_concat))) if len(_concat) else 0.1
ax.set_xlim(-20, 20)
ax.set_ylim(-_ymax * 1.15, _ymax * 1.15)
ax.set_xlabel("AoA − α₀ [deg]（零揚力角基準）", fontsize=TS)
ax.set_ylabel(r"$C_l$ − 過去平均", fontsize=TS)
ax.set_title("過去剛体翼の平均(=0)からの偏差：min/max帯 と 新システム", fontsize=TS)
ax.legend(fontsize=9, loc="upper left")
ax.grid(linestyle="--", alpha=0.4)
fig.tight_layout()
img_envelope = fig_to_bytes(fig)

# ─── PowerPoint 作成 ──────────────────────────────────────────────────────────
prs = Presentation(TEMPLATE_PATH)
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

sldIdLst = prs.slides._sldIdLst
for sId in list(sldIdLst):
    rId = sId.get(_qn('r:id'))
    prs.part.drop_rel(rId)
    sldIdLst.remove(sId)

LAY_TITLE   = prs.slide_layouts[0]
LAY_CONTENT = prs.slide_layouts[1]

CA_L = Inches(0.333)
CA_T = Inches(1.167)
CA_W = Inches(9.333)
CA_H = Inches(5.535)

DARK    = RGBColor(0x1A, 0x1A, 0x2E)
ACC     = RGBColor(0x4A, 0x9E, 0xEA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
TBLHDR  = RGBColor(0x1A, 0x3A, 0x6A)
NOTE_BG = RGBColor(0xF0, 0xF4, 0xFF)

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def _enable_shrink(tf):
    """テキストが箱を超えたら自動縮小（PowerPoint の『はみ出す場合縮小』）。
    データ増加で行数・件数が増えても箱からはみ出さないようにする保険。"""
    bodyPr = tf._txBody.find(_qn('a:bodyPr'))
    if bodyPr is None:
        return
    for _tag in ('a:normAutofit', 'a:spAutoFit', 'a:noAutofit'):
        _e = bodyPr.find(_qn(_tag))
        if _e is not None:
            bodyPr.remove(_e)
    bodyPr.append(bodyPr.makeelement(_qn('a:normAutofit'), {}))

def add_text(slide, text, l, t, w, h, size=14, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, wrap=True, shrink=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if shrink:
        _enable_shrink(tf)
    return txb

def add_img(slide, img_bytes, l, t, w, h):
    img_bytes.seek(0)
    slide.shapes.add_picture(img_bytes, l, t, w, h)

def _png_ar(img_bytes):
    """PNG の縦横比（幅/高さ）を IHDR から取得（PIL 不要）。"""
    img_bytes.seek(0)
    head = img_bytes.read(24)
    img_bytes.seek(0)
    w, h = struct.unpack(">II", head[16:24])
    return w / h

def add_img_fit(slide, img_bytes, l, t, max_w, max_h, halign="center", valign="center"):
    """アスペクト比を保ったまま (max_w, max_h) の枠に収め、枠内で中央寄せ配置。
    図を引き伸ばさずに自然に収める。"""
    ar  = _png_ar(img_bytes)
    box = max_w / max_h
    if ar >= box:
        w = int(max_w); h = int(max_w / ar)
    else:
        h = int(max_h); w = int(max_h * ar)
    lx = int(l) + (int(max_w) - w) // 2 if halign == "center" else int(l)
    ty = int(t) + (int(max_h) - h) // 2 if valign == "center" else int(t)
    img_bytes.seek(0)
    slide.shapes.add_picture(img_bytes, Emu(lx), Emu(ty), Emu(w), Emu(h))

def note_box(slide, text, l, t, w, h, size=12):
    add_rect(slide, l, t, w, h, fill=NOTE_BG, line=ACC, line_w=Pt(1))
    add_text(slide, text, l + Inches(0.1), t + Inches(0.1),
             w - Inches(0.2), h - Inches(0.2), size=size, color=DARK)

# ══════════════════════════════════════════════════════════════════
# スライド 1: タイトル
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_TITLE)
sl.placeholders[0].text = "NACA0012 Rigid Wing\n全データ比較（〜260610）／揚力傾斜の謎の解決"
sl.placeholders[1].text = "M1　岡本雄哉"
add_text(sl, ("Laboratory Seminar, Imamura Lab & Yamashita Lab\n"
              "The University of Tokyo, Engineering building 7,  June 5, 2026"),
         Inches(0.38), Inches(0.29), Inches(7.19), Inches(0.77),
         size=15, color=RGBColor(0x44, 0x44, 0x44))

# ══════════════════════════════════════════════════════════════════
# スライド 2: データセット一覧表
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_CONTENT)
sl.shapes.title.text = "データセット一覧（rigid 実験・自動更新）"

cols_w = [Inches(1.35), Inches(1.0), Inches(1.1), Inches(1.15),
          Inches(1.35), Inches(1.1), Inches(1.95)]
headers = ["データ", "日付", "Cl @ AoA=0", "零揚力迎角\n[deg]",
           "Clスロープ\n[/deg]", "P-M RMSE\n(Cl)", "備考"]
# 行数に応じて高さ・フォントを自動調整（件数が増えてもスライドに収める）
_n_rows  = len(stats) + 1                       # ヘッダ + データ行
row_h    = min(Inches(0.46), Emu(int(Inches(6.05) / _n_rows)))
_fs_hdr  = 10 if row_h >= Inches(0.36) else 8
_fs_cell = 10 if row_h >= Inches(0.36) else 8
top0     = CA_T

for ci, (hd, cw) in enumerate(zip(headers, cols_w)):
    lft = CA_L + sum(cols_w[:ci])
    add_rect(sl, lft, top0, cw, row_h, fill=TBLHDR,
             line=RGBColor(0x88, 0x99, 0xBB), line_w=Pt(0.5))
    add_text(sl, hd, lft + Inches(0.04), top0 + Inches(0.03),
             cw - Inches(0.08), row_h - Inches(0.06),
             size=_fs_hdr, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

stat_keys = list(stats.keys())
for ri, ((disp_name, who), k) in enumerate(zip(ROW_INFO, stat_keys)):
    s  = stats[k]
    bg = BG_COLOR[who]
    pm_str = f"{s['pm_rmse']:.3f}" if not np.isnan(s["pm_rmse"]) else "—"
    note   = NOTE_STR[who]
    vals   = [disp_name, s["date"],
              f"{s['Cl0']:+.4f}", f"{s['aoa0']:.3f}°",
              f"{s['slope']:.4f}", pm_str, note]
    for ci, (val, cw) in enumerate(zip(vals, cols_w)):
        lft = CA_L + sum(cols_w[:ci])
        top = top0 + row_h * (ri + 1)
        add_rect(sl, lft, top, cw, row_h, fill=bg,
                 line=RGBColor(0xCC, 0xCC, 0xCC), line_w=Pt(0.5))
        fsz = (_fs_cell - 1) if ci == 6 else _fs_cell
        add_text(sl, val, lft + Inches(0.04), top + Inches(0.04),
                 cw - Inches(0.08), row_h - Inches(0.08),
                 size=fsz, align=PP_ALIGN.CENTER)

# ── 統計値の取り出し（スライド本文で使用）──
a260605_1 = stats["260605①"]["aoa0"]
a260605_2 = stats["260605②"]["aoa0"]
a_r3      = stats["rigid3\n修正後"]["aoa0"]
sl260605_1 = stats["260605①"]["slope"]
sl260605_2 = stats["260605②"]["slope"]
sl_r3      = stats["rigid3\n修正後"]["slope"]
pct_r3     = (sl_r3 / OLD_AVG - 1.0) * 100.0

# ══════════════════════════════════════════════════════════════════
# スライド 3: Cl 曲線 全データ比較（メイン）
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_CONTENT)
sl.shapes.title.text = "Cl 曲線 全データ比較"
add_img_fit(sl, img_cl, CA_L, CA_T, Inches(6.5), Inches(5.2))
note_box(sl, ("【ポイント】\n"
              "・赤橙(接触不良)は傾きが過大\n\n"
              "・緑(260610_2 修正後)は\n"
              "　青(旧システム)とほぼ重なる\n\n"
              "→ 接触不良の差が修正され、\n"
              "　新システム＝旧システム相当に。\n\n"
              "・線形域の傾き・曲線形状とも\n"
              "　旧システムと一致"),
         CA_L + Inches(6.65), CA_T, Inches(2.6), Inches(4.5))

# ══════════════════════════════════════════════════════════════════
# スライド 4: 過去平均(=0)からの偏差：min/max帯 と 修正後新システム
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_CONTENT)
sl.shapes.title.text = "過去剛体翼の平均(=0)からの偏差 と 修正後新システム（260610_2）"
add_img_fit(sl, img_envelope, CA_L, CA_T, Inches(6.5), Inches(5.2))
note_box(sl, (f"【見方】過去平均=0 基準\n"
              f"（零揚力角基準 AoA−α₀）\n\n"
              f"・水色帯 = 過去剛体翼(旧)の\n"
              f"　min〜max（平均からの偏差）\n"
              f"・0 線 = 過去剛体翼の平均\n"
              f"・色付き線 = 新システム\n"
              f"　（修正後 260610_2 と今後の新規）\n\n"
              f"→ 修正後は過去平均(0)の近傍に収まり、\n"
              f"　接触不良時のような大きな差はない。\n"
              f"　（260610_2 の線形傾きは過去の\n"
              f"　 下端付近でわずかに小さめ＝{pct_r3:+.1f}%）"),
         CA_L + Inches(6.65), CA_T, Inches(2.6), Inches(5.0))

# ══════════════════════════════════════════════════════════════════
# スライド 5: 線形域拡大（Cl の傾き一致）
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_CONTENT)
sl.shapes.title.text = "AoA = −5〜+5° 拡大：線形域の Cl の傾き比較"
add_img_fit(sl, img_zoom, CA_L, CA_T, Inches(6.5), Inches(5.2))
note_box(sl, ("【読み方】\n"
              "線形域（±5°）の Cl の傾き\n\n"
              "・赤橙(接触不良)は\n"
              "　傾きが急（過大）\n\n"
              "・緑(260610_2 修正後)は\n"
              "　青(旧システム)とほぼ重なる\n\n"
              "→ 修正後は旧システムと\n"
              "　同じ傾きに戻った"),
         CA_L + Inches(6.65), CA_T, Inches(2.6), Inches(5.0))

# ══════════════════════════════════════════════════════════════════
# スライド 5: Cl スロープの時系列
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_CONTENT)
sl.shapes.title.text = "Cl スロープ（dCl/dAoA_mod）の時系列変化"
add_img_fit(sl, img_slope, CA_L, CA_T, Inches(6.5), Inches(4.5))
note_box(sl, (f"【揚力傾斜 dCl/dα】\n"
              f"・旧システム: 0.110〜0.112/deg\n"
              f"・接触不良(新): {sl260605_1:.3f}〜{sl260605_2:.3f}\n"
              f"　（差圧mV過小で過大に出た）\n"
              f"・修正後 260610_2: {sl_r3:.4f}\n\n"
              f"点線: 2π/rad = 0.1097/deg\n\n"
              f"→ 接触不良を修正すると、新システムの\n"
              f"　傾斜は旧システムとほぼ一致。\n"
              f"　新旧の差は解消された。"),
         CA_L + Inches(6.65), CA_T, Inches(2.6), Inches(4.5))

# ══════════════════════════════════════════════════════════════════
# スライド 6: 新システムの揚力傾斜 − 過去剛体翼の平均（新規）
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_CONTENT)
sl.shapes.title.text = "新システムの揚力傾斜 − 過去剛体翼の平均"
add_img_fit(sl, img_slope_diff, CA_L, CA_T, Inches(6.5), Inches(4.6))
note_box(sl, (f"【過去平均との差】\n"
              f"過去剛体翼(旧システム)の\n"
              f"平均傾斜 = {OLD_AVG:.4f}/deg\n"
              f"（±1σ = {OLD_STD:.4f}）\n\n"
              f"・接触不良(新)は +{sl260605_1-OLD_AVG:+.3f}〜\n"
              f"　と平均を大きく上回る\n\n"
              f"・修正後 260610_2: {sl_r3-OLD_AVG:+.4f}\n"
              f"　（{pct_r3:+.1f}%）= ほぼ過去平均\n\n"
              f"→ 修正後の新システムは過去剛体翼\n"
              f"　平均とほぼ一致（±1σ帯内）"),
         CA_L + Inches(6.65), CA_T, Inches(2.6), Inches(4.9))

# ══════════════════════════════════════════════════════════════════
# スライド 7: P-M 対称性誤差
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_CONTENT)
sl.shapes.title.text = "正負掃引の非対称性（P-M 対称性誤差）"
GW, GH = Inches(4.55), Inches(3.5)
add_img_fit(sl, img_pm,      CA_L,                CA_T, GW, GH)
add_img_fit(sl, img_pm_corr, CA_L + Inches(4.75), CA_T, GW, GH)
note_box(sl, ("→ 零揚力迎角を基準に補正すると（b）、取り付け角オフセット α₀ による\n"
              "　見かけ上の非対称性が除去され、実質的な空力非対称性を評価できる。\n"
              "→ 修正後の新システム(260610_2)の P-M RMSE は旧システムと同程度。新旧の差は解消し信頼性は回復。"),
         CA_L, CA_T + Inches(3.65), CA_W, Inches(1.75), size=12)

# ══════════════════════════════════════════════════════════════════
# スライド 9: 零揚力迎角（補足）
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_CONTENT)
sl.shapes.title.text = "（補足）零揚力迎角 α₀ について"
add_img_fit(sl, img_aoa0, CA_L, CA_T, Inches(6.5), Inches(4.5))
note_box(sl, (f"【補足：α₀ の見かけのずれ】\n"
              f"・0605 の実験のみ AoA=0 の\n"
              f"　原点パルスを 11250 に設定\n"
              f"・その他は全て 11025\n"
              f"　（差 225 pulse ÷ 250 = 0.9°）\n\n"
              f"→ 0605 の α₀({a260605_1:.2f}/{a260605_2:.2f}°)が\n"
              f"　大きいのは主にこの原点設定差。\n"
              f"　0.9°差し引くと旧域(0.5〜0.8°)。\n\n"
              f"※ 原点を 11025 に統一すれば\n"
              f"　α₀ も新旧でほぼ一致。"),
         CA_L + Inches(6.65), CA_T, Inches(2.6), Inches(4.7))

# ══════════════════════════════════════════════════════════════════
# スライド 11: まとめ
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_CONTENT)
sl.shapes.title.text = "まとめ"

items = [
    ("① 新旧の差の正体 ＝ 差圧デジボルの接触不良",
     f"接触不良時(260605_1/_2・260608・260610_1)は Cl 傾斜が旧システム(0.110〜0.112/deg)より大幅に過大({sl260605_1:.3f}/deg 等)。\n"
     "原因：差圧mVが接触不良で過小 → 風速U・動圧qが過小 → Cl=L/q が過大に。新旧の見かけの差はこれが正体だった。"),
    ("② 接触を修正すると新システムは旧システムとほぼ一致",
     f"接触を修正した 260610_2（真mV≈1320, U≈13.1m/s）の傾斜 = {sl_r3:.4f}/deg ≒ 旧システム(0.110〜0.112)・2π。\n"
     "→ 接触不良の差が修正され、新システムでも旧システムとほぼ同じ空力データが得られた。"),
    ("③ 零揚力迎角 α₀ の差は原点パルス設定の違い（補足）",
     "0605 のみ AoA=0 の原点を 11250、その他は全て 11025（差 225pulse = 0.9°）。\n"
     f"0605 の α₀({a260605_1:.2f}/{a260605_2:.2f}°)が大きいのは主にこの設定差で、0.9°差し引けば旧域。原点統一で新旧一致。"),
    ("④ 結論・今後",
     "新システムは接触不良さえ除けば旧システムと同等の空力データを再現。計測前に『差圧mVの妥当性確認』を追加。\n"
     "マウント弾性/センサ/定格/後処理は全て無罪と確認済み。過去の接触不良データも真mVが分かれば補正可能。"),
]

ITEM_H   = Inches(1.2)
ITEM_GAP = Inches(0.06)
for i, (title, body) in enumerate(items):
    top = CA_T + (ITEM_H + ITEM_GAP) * i
    add_rect(sl, CA_L, top, CA_W, ITEM_H,
             fill=RGBColor(0xE8, 0xF0, 0xFE), line=ACC, line_w=Pt(1))
    add_rect(sl, CA_L, top, Inches(0.07), ITEM_H, fill=ACC)
    add_text(sl, title, CA_L + Inches(0.15), top + Inches(0.04),
             CA_W - Inches(0.2), Inches(0.32), size=13, bold=True, color=DARK)
    add_text(sl, body,  CA_L + Inches(0.15), top + Inches(0.38),
             CA_W - Inches(0.2), Inches(0.76), size=10, color=DARK)

# ─── 保存 ─────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"保存完了: {OUT}")
